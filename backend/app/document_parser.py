from __future__ import annotations

import hashlib
import logging
import re
import shutil
import subprocess
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from .pdf_math_geometry import extract_math_layout_candidates
from .schemas import ParsedDocument, SourceBlock


TEXT_TYPES = {".txt", ".md", ".markdown"}
SUPPORTED_TYPES = TEXT_TYPES | {".pdf", ".pptx", ".docx"}
_LIST_PREFIX = re.compile(
    r"^(?:[-•●▪◦‣⁃]|(?:\d+|[一二三四五六七八九十]+)[、.)）])\s*"
)
_LATIN_HYPHENATED_END = re.compile(r"[A-Za-z]-$")
_LATIN_START = re.compile(r"^[a-z]")
_CJK_END = re.compile(r"[\u3400-\u9fff]$")
_CJK_START = re.compile(r"^[\u3400-\u9fff]")
_SENTENCE_END = re.compile(r"[。！？!?；;.]$")
_STANDALONE_PAGE_NUMBER = re.compile(
    r"^(?:第\s*)?\d{1,4}(?:\s*页)?$",
    re.IGNORECASE,
)
_CHAPTER_HEADING = re.compile(
    r"^第\s*(?P<number>\d{1,3}|[一二三四五六七八九十百]+)\s*章"
    r"(?:\s*[：:—-]?\s*(?P<title>[^。！？!?]{1,48}))?$"
)
_SECTION_HEADING = re.compile(
    r"^§\s*(?P<number>\d{1,3}(?:\s*\.\s*\d{1,3})+)"
    r"(?:\s*[：:—-]?\s*(?P<title>[^。！？!?]{1,48}))?$"
)
_HEADING_DECORATION = re.compile(r"^[*△▲◆◇■□●•]+\s*")
_SEMANTIC_WORD = re.compile(
    r"[\u3400-\u9fff]{2,}|[A-Za-z]{3,}|[Α-ω]{2,}"
)
_EQUATION_RELATION = re.compile(r"[=≈≃≅≤≥<>]")
_PRIVATE_USE_OR_REPLACEMENT = re.compile(r"[\ue000-\uf8ff\ufffd]")
_MATH_LAYOUT_CUE = re.compile(
    r"(?:"
    r"公式|方程|等于|关系|光程|线宽|频率|稳频|波长|能量|概率|"
    r"量子|积分|归一化|本征|角动量|磁矩|能级|跃迁|干涉|"
    r"[=≈≃≅≤≥<>∫∑Σ√Δλν]"
    r")",
    re.IGNORECASE,
)
_SYMBOL_FONT_PUA_TRANSLATION = str.maketrans(
    {
        "\uf03e": ">",
        "\uf06c": "λ",
        "\uf06e": "ν",
        "\uf040": "≈",
        "\uf044": "Δ",
        "\uf057": "Ω",
        "\uf0b4": "×",
        "\uf0bb": "≈",
        "\uf0d7": "·",
        "\uf0de": "⇒",
    }
)
_SCIENTIFIC_EXPONENT_CONTEXT = re.compile(r"数量级|小到|低至")
_SPACED_NEGATIVE_POWER = re.compile(
    r"(?<![\d.])10\s*[-−]\s*(?P<exponent>\d{1,3})(?!\d)"
)
_INLINE_SCIENTIFIC_POWER = re.compile(
    r"(?P<mantissa>\d+(?:\.\d+)?)"
    r"(?P<operator>\s*[×x*]\s*)"
    r"10(?P<exponent>\d{1,2})"
    r"(?=(?:\s|$|[A-Za-z\u3400-\u9fff]))"
)
_COLLAPSED_SCIENTIFIC_UNIT_POWER = re.compile(
    r"(?<!\d)10(?P<exponent>\d{1,2})"
    r"\s*W/(?P<unit>c?m)2(?!\d)",
    re.IGNORECASE,
)
_STACKED_DELTA_NUMERATOR = re.compile(
    r"^(?P<indent>[ \t]*)Δ(?P<symbol>[A-Za-zΑ-ω])(?=\s|\d|$)"
)
_STACKED_FRACTION_AFTER_RELATION = re.compile(
    r"(?m)"
    r"^[ \t]*(?P<numerator>[λν]\s*[A-Za-z0-9])"
    r"(?:\s*\([^\n]*\))?\s*$\n"
    r"^[^\n]*?(?P<body>(?<![A-Za-z0-9])"
    r"[A-Za-z][A-Za-z0-9]*\s*=\s*[A-Za-z0-9]+)\s*$\n"
    r"^[ \t]*(?P<denominator>[A-Za-z0-9]+)\s*$"
)
_STACKED_FRACTION_BEFORE_RELATION = re.compile(
    r"(?m)"
    r"^[ \t]*(?P<numerator>"
    r"[0-9A-Za-zΑ-ω][0-9A-Za-zΑ-ω \t]{0,20})\s*$\n"
    r"^[ \t]*(?P<lhs>[λν]\s*[A-Za-z0-9])\s*=\s*$\n"
    r"^[ \t]*(?P<denominator>[A-Za-z0-9]+)"
    r"(?:[ \t]{2,}[^\n]*)?$"
)
_PDFTOTEXT_TIMEOUT_SECONDS = 20
_PDFTOTEXT_MAX_OUTPUT_CHARS = 2_000_000

logger = logging.getLogger(__name__)


def _document_id(path: Path) -> str:
    digest = hashlib.sha1(path.read_bytes()).hexdigest()[:12]
    return f"doc_{digest}"


def _reflow_pdf_text(text: str) -> str:
    """Repair common extraction wraps without inventing cross-page structure."""

    logical_lines: list[str] = []
    current = ""

    def flush() -> None:
        nonlocal current
        if current:
            logical_lines.append(current.strip())
            current = ""

    for raw_line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        line = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", raw_line)
        line = re.sub(r"[ \t]+", " ", line).strip()
        if not line:
            flush()
            if logical_lines and logical_lines[-1] != "":
                logical_lines.append("")
            continue
        if _LIST_PREFIX.match(line):
            flush()
            logical_lines.append(line)
            continue
        if not current:
            current = line
            continue
        if _SENTENCE_END.search(current):
            flush()
            current = line
            continue
        if _LATIN_HYPHENATED_END.search(current) and _LATIN_START.search(line):
            current = f"{current[:-1]}{line}"
            continue
        separator = "" if _CJK_END.search(current) and _CJK_START.search(line) else " "
        current = f"{current}{separator}{line}"

    flush()
    while logical_lines and logical_lines[-1] == "":
        logical_lines.pop()
    return "\n".join(logical_lines).strip()


def _without_standalone_page_numbers(text: str) -> str:
    lines = [
        line.strip()
        for line in text.replace("\r\n", "\n").replace("\r", "\n").splitlines()
        if line.strip() and not _STANDALONE_PAGE_NUMBER.fullmatch(line.strip())
    ]
    return "\n".join(lines).strip()


def _looks_like_meaningful_equation(text: str) -> bool:
    compact = re.sub(r"\s+", "", _without_standalone_page_numbers(text))
    if not compact:
        return False
    if _EQUATION_RELATION.search(compact):
        left, right = _EQUATION_RELATION.split(compact, maxsplit=1)
        return bool(
            re.search(r"[\w\u3400-\u9fffΑ-ω]", left)
            and re.search(r"[\w\u3400-\u9fffΑ-ω]", right)
        )
    return bool(
        re.search(r"[∫∑Σ√]", compact)
        and len(re.findall(r"[\w\u3400-\u9fffΑ-ω]", compact)) >= 3
    )


def _looks_like_axis_or_symbol_fragment(text: str) -> bool:
    body = _without_standalone_page_numbers(text)
    if not body or _SEMANTIC_WORD.search(body):
        return False
    if _looks_like_meaningful_equation(body):
        return False
    tokens = [
        token.casefold()
        for token in re.findall(
            r"[A-Za-z]+|[Α-ω]+|[\u3400-\u9fff]+",
            body,
        )
    ]
    if not tokens:
        return bool(re.search(r"[^\W\d_]", body, re.UNICODE)) or bool(
            re.search(r"[^\w\s]", body)
        )
    axis_tokens = {
        "x",
        "y",
        "z",
        "r",
        "t",
        "o",
        "θ",
        "φ",
        "ϕ",
        "ρ",
        "ψ",
    }
    return (
        len(tokens) >= 2
        and all(token in axis_tokens for token in tokens)
        and (
            len(tokens) <= 3
            or len(set(tokens)) <= max(2, len(tokens) // 2)
        )
    )


def _is_clearly_low_quality_pdf_text(text: str) -> bool:
    """Gate the expensive secondary extractor to unmistakably weak pages."""

    body = _without_standalone_page_numbers(text)
    compact = re.sub(r"\s+", "", body)
    if not compact:
        return True
    if _looks_like_meaningful_equation(body):
        return False
    if not re.search(r"[\w\u3400-\u9fffΑ-ω]", compact, re.UNICODE):
        return True
    if _looks_like_axis_or_symbol_fragment(body):
        return True

    private_use_count = len(re.findall(r"[\ue000-\uf8ff\ufffd]", compact))
    semantic_count = len(
        re.findall(r"[\w\u3400-\u9fffΑ-ω]", compact, re.UNICODE)
    )
    if (
        private_use_count / max(len(compact), 1) >= 0.35
        and semantic_count / max(len(compact), 1) < 0.45
    ):
        return True
    return False


def _needs_math_layout_comparison(text: str) -> bool:
    """Compare layout extractors when embedded math glyphs may be dropped.

    A page can contain plenty of readable prose and still lose the one
    equation or ratio that carries the claim.  Private-use glyphs plus
    numerical/math context are a cheap signal for that failure mode.
    """

    body = _without_standalone_page_numbers(text)
    if not body or not _PRIVATE_USE_OR_REPLACEMENT.search(body):
        return False
    return bool(re.search(r"\d", body) and _MATH_LAYOUT_CUE.search(body))


def _compact_math_fragment(text: str) -> str:
    compact = re.sub(r"\s+", "", text)
    return re.sub(r"([λν])([A-Za-z0-9])", r"\1_\2", compact)


def _stacked_delta_ratio_hints(text: str) -> list[str]:
    lines = text.splitlines()
    hints: list[str] = []
    for index, line in enumerate(lines):
        numerator = _STACKED_DELTA_NUMERATOR.match(line)
        if not numerator:
            continue
        symbol = numerator.group("symbol")
        numerator_indent = len(numerator.group("indent").expandtabs(4))
        denominator_pattern = re.compile(
            rf"^(?P<indent>[ \t]*){re.escape(symbol)}(?=\s|$)"
        )
        for candidate in lines[index + 1 : index + 4]:
            denominator = denominator_pattern.match(candidate)
            if not denominator:
                continue
            denominator_indent = len(
                denominator.group("indent").expandtabs(4)
            )
            if abs(numerator_indent - denominator_indent) <= 3:
                hint = f"Δ{symbol}/{symbol}"
                if hint not in text and hint not in hints:
                    hints.append(hint)
                break
    return hints


def _stacked_fraction_hints(text: str) -> list[str]:
    hints: list[str] = []
    for match in _STACKED_FRACTION_AFTER_RELATION.finditer(text):
        body = re.sub(r"\s*=\s*", " = ", match.group("body"))
        numerator = _compact_math_fragment(match.group("numerator"))
        denominator = _compact_math_fragment(match.group("denominator"))
        hint = f"{body}{numerator}/{denominator}"
        if hint not in text and hint not in hints:
            hints.append(hint)

    for match in _STACKED_FRACTION_BEFORE_RELATION.finditer(text):
        lhs = _compact_math_fragment(match.group("lhs"))
        numerator = _compact_math_fragment(match.group("numerator"))
        denominator = _compact_math_fragment(match.group("denominator"))
        hint = f"{lhs} = {numerator}/{denominator}"
        if hint not in text and hint not in hints:
            hints.append(hint)
    return hints


def _prepare_pdf_math_layout(text: str) -> str:
    """Normalize known Symbol-font glyphs and preserve stacked equations.

    Poppler keeps the geometry of legacy courseware equations but can emit
    Symbol-font characters in Unicode's private-use area.  Reflowing that
    output directly destroys stacked fractions, so add a conservative
    one-line equivalent only when numerator/denominator alignment is clear.
    """

    normalized = text.translate(_SYMBOL_FONT_PUA_TRANSLATION)
    lines: list[str] = []
    for line in normalized.splitlines():
        line = _INLINE_SCIENTIFIC_POWER.sub(
            lambda match: (
                f"{match.group('mantissa')}×10^"
                f"{match.group('exponent')}"
            ),
            line,
        )
        line = _COLLAPSED_SCIENTIFIC_UNIT_POWER.sub(
            lambda match: (
                f"10^{match.group('exponent')} "
                f"W/{match.group('unit')}²"
            ),
            line,
        )
        if _SCIENTIFIC_EXPONENT_CONTEXT.search(line):
            line = _SPACED_NEGATIVE_POWER.sub(
                lambda match: f"10^-{match.group('exponent')}",
                line,
            )
        lines.append(line)
    normalized = "\n".join(lines)

    hints = [
        *_stacked_delta_ratio_hints(normalized),
        *_stacked_fraction_hints(normalized),
    ]
    if hints:
        normalized = f"{normalized.rstrip()}\n\n" + "；".join(hints)
    return normalized


def _math_layout_score(text: str) -> int:
    """Score recoverable mathematical payload, not general text length."""

    body = _without_standalone_page_numbers(_prepare_pdf_math_layout(text))
    compact = re.sub(r"\s+", "", body)
    relation_count = len(_EQUATION_RELATION.findall(compact))
    fraction_count = compact.count("/")
    math_symbol_count = len(
        re.findall(r"[∫∑Σ√Δλνπħ×÷^]", compact, re.IGNORECASE)
    )
    scientific_count = len(
        re.findall(
            r"\d(?:\.\d+)?\s*(?:×|x|\*)\s*10\s*(?:\^|[-−])?\s*\d+",
            body,
            re.IGNORECASE,
        )
    )
    return (
        relation_count * 8
        + fraction_count * 4
        + scientific_count * 3
        + min(math_symbol_count, 12)
    )


def _fallback_improves_math_layout(primary: str, fallback: str) -> bool:
    if _is_clearly_low_quality_pdf_text(fallback):
        return False
    return _math_layout_score(fallback) > _math_layout_score(primary)


def _infer_pdf_heading(text: str) -> str | None:
    lines = [
        line.strip()
        for line in text.replace("\r\n", "\n").replace("\r", "\n").splitlines()
        if line.strip()
    ]
    first_line = lines[0] if lines else ""
    if not first_line:
        return None
    first_line = _HEADING_DECORATION.sub("", first_line)
    first_line = re.sub(r"\s+", " ", first_line).strip()
    if not first_line or len(first_line) > 72 or "结束" in first_line:
        return None
    if first_line.endswith(("，", "、", ",", "：", ":")):
        continuation = (
            re.sub(r"\s+", " ", lines[1]).strip()
            if len(lines) > 1
            else ""
        )
        safe_continuation = bool(
            continuation
            and len(continuation) <= 32
            and _SEMANTIC_WORD.search(continuation)
            and not re.match(
                r"^(?:§|第\s*\d+\s*章|\d+[、.)）])",
                continuation,
            )
            and not re.search(r"[。！？!?；;:]$", continuation)
        )
        if safe_continuation:
            separator = "" if first_line.endswith(("，", "、", ",")) else " "
            first_line = f"{first_line}{separator}{continuation}"
        else:
            return None

    chapter = _CHAPTER_HEADING.fullmatch(first_line)
    if chapter:
        marker = f"第{chapter.group('number')}章"
        title = (chapter.group("title") or "").strip()
        return f"{marker} {title}".strip()

    section = _SECTION_HEADING.fullmatch(first_line)
    if section:
        number = re.sub(r"\s+", "", section.group("number"))
        marker = f"§{number}"
        title = (section.group("title") or "").strip()
        return f"{marker} {title}".strip()
    return None


def _pdftotext_page(
    executable: str,
    path: Path,
    page_number: int,
) -> str:
    completed = subprocess.run(
        [
            executable,
            "-f",
            str(page_number),
            "-l",
            str(page_number),
            "-layout",
            "-enc",
            "UTF-8",
            str(path),
            "-",
        ],
        check=True,
        timeout=_PDFTOTEXT_TIMEOUT_SECONDS,
        capture_output=True,
        stdin=subprocess.DEVNULL,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    if len(completed.stdout) > _PDFTOTEXT_MAX_OUTPUT_CHARS:
        raise ValueError("pdftotext 单页输出超过安全上限")
    return completed.stdout


def _record_parse_warning(warnings: list[str], warning: str) -> None:
    if warning not in warnings:
        warnings.append(warning)
        logger.warning("%s", warning)


def _parse_pdf(
    path: Path,
    *,
    parse_metadata: dict[str, Any] | None = None,
    warnings: list[str] | None = None,
) -> list[SourceBlock]:
    reader = PdfReader(str(path))
    metadata = parse_metadata if parse_metadata is not None else {}
    parse_warnings = warnings if warnings is not None else []
    fallback_metadata: dict[str, Any] = {
        "tool": "pdftotext-layout",
        "tool_available": None,
        "candidate_pages": [],
        "low_quality_pages": [],
        "math_comparison_pages": [],
        "attempted_pages": [],
        "used_pages": [],
        "failed_pages": [],
        "retained_low_quality_pages": [],
        "retained_primary_pages": [],
    }
    geometry_metadata: dict[str, Any] = {
        "tool": "pdfplumber-geometry",
        "tool_available": None,
        "tool_version": None,
        "candidate_pages": [],
        "attempted_pages": [],
        "used_pages": [],
        "failed_pages": [],
        "candidate_count": 0,
        "candidates": [],
        "injected_into_text": False,
    }
    metadata.update(
        {
            "pdf_page_count": len(reader.pages),
            "pdf_primary_extractor": "pypdf-layout",
            "pdf_pypdf_failed_pages": [],
            "pdf_text_fallback": fallback_metadata,
            "pdf_geometry_math": geometry_metadata,
        }
    )
    blocks: list[SourceBlock] = []
    current_heading: str | None = None
    pdftotext_checked = False
    pdftotext: str | None = None
    geometry_checked = False
    geometry_document: Any | None = None
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            extracted = page.extract_text(extraction_mode="layout")
        except TypeError:
            # Compatibility with older pypdf graph versions/providers.
            try:
                extracted = page.extract_text()
            except Exception as exc:  # pragma: no cover - provider-specific
                extracted = ""
                metadata["pdf_pypdf_failed_pages"].append(page_number)
                _record_parse_warning(
                    parse_warnings,
                    "[pdf_parse_degraded:pypdf_failed] "
                    f"第 {page_number} 页 pypdf 提取失败"
                    f"（{type(exc).__name__}），尝试安全回退。",
                )
        except Exception as exc:  # pragma: no cover - provider-specific
            extracted = ""
            metadata["pdf_pypdf_failed_pages"].append(page_number)
            _record_parse_warning(
                parse_warnings,
                "[pdf_parse_degraded:pypdf_failed] "
                f"第 {page_number} 页 pypdf 提取失败"
                f"（{type(exc).__name__}），尝试安全回退。",
            )

        selected_raw = extracted or ""
        primary_low_quality = _is_clearly_low_quality_pdf_text(
            selected_raw
        )
        math_comparison = _needs_math_layout_comparison(selected_raw)
        if primary_low_quality or math_comparison:
            fallback_metadata["candidate_pages"].append(page_number)
            if primary_low_quality:
                fallback_metadata["low_quality_pages"].append(page_number)
            if math_comparison:
                fallback_metadata["math_comparison_pages"].append(
                    page_number
                )
            if not pdftotext_checked:
                pdftotext = shutil.which("pdftotext")
                pdftotext_checked = True
                fallback_metadata["tool_available"] = bool(pdftotext)
            if pdftotext:
                fallback_metadata["attempted_pages"].append(page_number)
                try:
                    fallback_raw = _pdftotext_page(
                        pdftotext,
                        path,
                        page_number,
                    )
                except (
                    OSError,
                    subprocess.SubprocessError,
                    ValueError,
                ) as exc:
                    fallback_metadata["failed_pages"].append(page_number)
                    _record_parse_warning(
                        parse_warnings,
                        "[pdf_parse_degraded:pdftotext_failed] "
                        f"第 {page_number} 页 pdftotext -layout 回退失败"
                        f"（{type(exc).__name__}），已保留 pypdf 结果。",
                    )
                else:
                    fallback_usable = bool(
                        fallback_raw.strip()
                        and not _is_clearly_low_quality_pdf_text(
                            fallback_raw
                        )
                    )
                    use_fallback = bool(
                        fallback_usable
                        and (
                            primary_low_quality
                            or (
                                math_comparison
                                and _fallback_improves_math_layout(
                                    selected_raw,
                                    fallback_raw,
                                )
                            )
                        )
                    )
                    if use_fallback:
                        selected_raw = fallback_raw
                        fallback_metadata["used_pages"].append(page_number)
                    elif primary_low_quality:
                        fallback_metadata[
                            "retained_low_quality_pages"
                        ].append(page_number)
                        _record_parse_warning(
                            parse_warnings,
                            "[pdf_parse_degraded:low_quality_retained] "
                            f"第 {page_number} 页两种文本提取结果均明显低质，"
                            "已保留 pypdf 结果并等待视觉单元覆盖。",
                        )
                    else:
                        fallback_metadata["retained_primary_pages"].append(
                            page_number
                        )
            else:
                if primary_low_quality:
                    fallback_metadata["retained_low_quality_pages"].append(
                        page_number
                    )
                else:
                    fallback_metadata["retained_primary_pages"].append(
                        page_number
                    )

        selected_raw = _prepare_pdf_math_layout(selected_raw)
        if math_comparison:
            geometry_metadata["candidate_pages"].append(page_number)
            if not geometry_checked:
                geometry_checked = True
                geometry_error: Exception | None = None
                try:
                    import pdfplumber
                except ImportError as exc:
                    geometry_error = exc
                else:
                    try:
                        from pdfplumber.utils.exceptions import (
                            PdfminerException,
                        )

                        geometry_document = pdfplumber.open(path)
                    except (
                        OSError,
                        ValueError,
                        PdfminerException,
                    ) as exc:
                        geometry_error = exc
                    else:
                        geometry_metadata["tool_available"] = True
                        geometry_metadata["tool_version"] = getattr(
                            pdfplumber,
                            "__version__",
                            None,
                        )
                if geometry_error is not None:
                    geometry_metadata["tool_available"] = False
                    _record_parse_warning(
                        parse_warnings,
                        "[pdf_parse_degraded:geometry_unavailable] "
                        "字符几何恢复层不可用"
                        f"（{type(geometry_error).__name__}），数学公式需要复核。",
                    )
            if geometry_document is not None:
                geometry_metadata["attempted_pages"].append(page_number)
                try:
                    geometry_candidates = extract_math_layout_candidates(
                        geometry_document.pages[page_number - 1]
                    )
                except (IndexError, OSError, TypeError, ValueError) as exc:
                    geometry_metadata["failed_pages"].append(page_number)
                    _record_parse_warning(
                        parse_warnings,
                        "[pdf_parse_degraded:geometry_failed] "
                        f"第 {page_number} 页字符几何恢复失败"
                        f"（{type(exc).__name__}），已保留文本抽取结果。",
                    )
                else:
                    if geometry_candidates:
                        geometry_metadata["used_pages"].append(page_number)
                        geometry_metadata["candidate_count"] += len(
                            geometry_candidates
                        )
                        geometry_metadata["candidates"].extend(
                            {
                                "page": page_number,
                                "canonical": candidate.canonical,
                                "source_bbox": [
                                    round(value, 3)
                                    for value in candidate.source_bbox
                                ],
                                "confidence": candidate.confidence,
                                "kind": candidate.kind,
                                "origin": candidate.origin,
                                "issues": list(candidate.issues),
                            }
                            for candidate in geometry_candidates
                        )
        detected_heading = _infer_pdf_heading(selected_raw)
        if detected_heading:
            current_heading = detected_heading
        text = _reflow_pdf_text(selected_raw)
        if text:
            blocks.append(
                SourceBlock(
                    text=text,
                    page=page_number,
                    heading=current_heading,
                )
            )

    unavailable_low_quality_pages = [
        page
        for page in fallback_metadata["low_quality_pages"]
        if page not in fallback_metadata["attempted_pages"]
    ]
    if unavailable_low_quality_pages:
        pages = "、".join(
            str(page) for page in unavailable_low_quality_pages
        )
        _record_parse_warning(
            parse_warnings,
            "[pdf_parse_degraded:pdftotext_unavailable] "
            f"第 {pages} 页 pypdf 文本明显低质，但未找到 pdftotext；"
            "已保留原提取结果并等待视觉单元覆盖。",
        )
    unavailable_math_pages = [
        page
        for page in fallback_metadata["math_comparison_pages"]
        if page not in fallback_metadata["attempted_pages"]
        and page not in unavailable_low_quality_pages
    ]
    if unavailable_math_pages:
        pages = "、".join(str(page) for page in unavailable_math_pages)
        _record_parse_warning(
            parse_warnings,
            "[pdf_parse_degraded:pdftotext_unavailable_math] "
            f"第 {pages} 页含嵌入式数学字形，但未找到 pdftotext；"
            "已保留 pypdf 结果，公式与比值需要复核。",
        )
    if geometry_document is not None:
        geometry_document.close()
    return blocks


def _shape_text(shape) -> str:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            parts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                parts.append(" | ".join(values))
    return "\n".join(parts).strip()


def _parse_pptx(path: Path) -> list[SourceBlock]:
    presentation = Presentation(str(path))
    blocks: list[SourceBlock] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title: str | None = None
        title_shape_id: int | None = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip() or None
            title_shape_id = slide.shapes.title.shape_id

        ordered: list[tuple[int, int, int, int, str]] = []
        for shape_index, shape in enumerate(slide.shapes):
            text = _shape_text(shape)
            if not text:
                continue
            is_title = int(getattr(shape, "shape_id", None) != title_shape_id)
            top = int(getattr(shape, "top", 0) or 0)
            left = int(getattr(shape, "left", 0) or 0)
            ordered.append((is_title, top, left, shape_index, text))

        texts: list[str] = []
        seen: set[str] = set()
        for _, _, _, _, text in sorted(ordered):
            normalized = re.sub(r"\s+", " ", text).casefold()
            if normalized in seen:
                continue
            seen.add(normalized)
            texts.append(text)
        combined = "\n".join(texts).strip()
        if combined:
            blocks.append(
                SourceBlock(text=combined, slide=slide_number, heading=title)
            )
    return blocks


def _parse_docx(path: Path) -> list[SourceBlock]:
    document = Document(str(path))
    blocks: list[SourceBlock] = []
    current_heading: str | None = None
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style and paragraph.style.name.startswith("Heading"):
            current_heading = text
        blocks.append(SourceBlock(text=text, heading=current_heading))
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                blocks.append(
                    SourceBlock(
                        text=" | ".join(values),
                        heading=current_heading,
                    )
                )
    return blocks


def _parse_text(path: Path) -> list[SourceBlock]:
    text = path.read_text(encoding="utf-8", errors="replace")
    blocks: list[SourceBlock] = []
    current_heading: str | None = None
    for raw_block in text.replace("\r\n", "\n").split("\n\n"):
        block = raw_block.strip()
        if not block:
            continue
        first_line = block.splitlines()[0].strip()
        if first_line.startswith("#"):
            current_heading = first_line.lstrip("#").strip()
        blocks.append(SourceBlock(text=block, heading=current_heading))
    return blocks


def parse_visual_document(
    path: Path,
    original_filename: str | None = None,
) -> ParsedDocument:
    """Create the PDF document shell without extracting page text."""

    suffix = path.suffix.lower()
    if suffix != ".pdf":
        raise ValueError("视觉直抽文档壳仅支持 PDF。")

    reader = PdfReader(str(path))
    page_count = len(reader.pages)
    if page_count < 1:
        raise ValueError("PDF 没有可渲染页面。")

    filename = original_filename or path.name
    return ParsedDocument(
        document_id=_document_id(path),
        filename=filename,
        file_type="pdf",
        title=Path(filename).stem,
        blocks=[],
        parse_metadata={
            "pdf_page_count": page_count,
            "pdf_input_mode": "direct_visual_only",
            "pdf_text_extraction_performed": False,
        },
    )


def parse_document(path: Path, original_filename: str | None = None) -> ParsedDocument:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_TYPES:
        raise ValueError(
            f"暂不支持 {suffix or '未知'} 文件，请上传 PDF、PPTX、DOCX、TXT 或 MD。"
        )

    parse_metadata: dict[str, Any] = {}
    warnings: list[str] = []
    if suffix == ".pdf":
        blocks = _parse_pdf(
            path,
            parse_metadata=parse_metadata,
            warnings=warnings,
        )
    elif suffix == ".pptx":
        blocks = _parse_pptx(path)
    elif suffix == ".docx":
        blocks = _parse_docx(path)
    else:
        blocks = _parse_text(path)

    if not blocks:
        raise ValueError("没有从文档中解析出可用文本。扫描版 PDF 需要先接入 OCR。")

    filename = original_filename or path.name
    title = next((block.heading for block in blocks if block.heading), None)
    if not title:
        title = Path(filename).stem

    return ParsedDocument(
        document_id=_document_id(path),
        filename=filename,
        file_type=suffix.lstrip("."),
        title=title,
        blocks=blocks,
        parse_metadata=parse_metadata,
        warnings=warnings,
    )

def parse_documents(
    paths: list[Path],
    original_filenames: list[str] | None = None,
) -> ParsedDocument:
    if not paths:
        raise ValueError("未提供任何待解析文档。")
    if len(paths) == 1:
        return parse_document(
            paths[0],
            original_filenames[0] if original_filenames else None,
        )

    all_blocks: list[SourceBlock] = []
    all_warnings: list[str] = []
    combined_metadata: dict[str, Any] = {"multi_document": True, "documents": []}
    filenames = original_filenames or [p.name for p in paths]

    for idx, (p, fn) in enumerate(zip(paths, filenames, strict=False)):
        doc = parse_document(p, fn)
        all_warnings.extend(doc.warnings)
        combined_metadata["documents"].append(
            {
                "index": idx + 1,
                "filename": fn,
                "document_id": doc.document_id,
                "block_count": len(doc.blocks),
            }
        )
        for b in doc.blocks:
            heading = f"[{fn}] {b.heading}" if b.heading else f"[{fn}]"
            all_blocks.append(
                SourceBlock(
                    heading=heading,
                    text=b.text,
                    page=b.page,
                    slide=b.slide,
                )
            )

    titles = [fn for fn in filenames[:3]]
    title_summary = " & ".join(titles) + (f" 等共 {len(paths)} 份文档" if len(paths) > 3 else "")
    combined_id = hashlib.sha256("::".join(str(p) for p in paths).encode()).hexdigest()[:16]

    return ParsedDocument(
        document_id=combined_id,
        filename=title_summary,
        file_type="multi",
        title=title_summary,
        blocks=all_blocks,
        parse_metadata=combined_metadata,
        warnings=all_warnings,
    )
