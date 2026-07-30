from __future__ import annotations

import hashlib
import json
import math
import os
import shutil
import statistics
import subprocess
import uuid
from pathlib import Path
from urllib.parse import quote, urlsplit, urlunsplit

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

from .schemas import (
    CropRequest,
    RenderResponse,
    RenderedPage,
    VisualAsset,
    VisualUnit,
)


IMAGE_TYPES = {".png", ".jpg", ".jpeg", ".webp"}
RENDER_TYPES = IMAGE_TYPES | {".pdf", ".pptx"}
VISUAL_DEGRADED_RENDER_BUDGET = "[visual_degraded:render_budget]"
VISUAL_DEGRADED_RENDER_FAILURE = "[visual_degraded:render_failure]"
VISUAL_DEGRADED_NATIVE_EXTRACTION = (
    "[visual_degraded:native_extraction]"
)


def _select_page_numbers(total_pages: int, max_pages: int) -> list[int]:
    if total_pages <= 0 or max_pages <= 0:
        return []
    if total_pages <= max_pages:
        return list(range(1, total_pages + 1))
    if max_pages == 1:
        return [1]
    last_index = total_pages - 1
    indices = [
        round(position * last_index / (max_pages - 1))
        for position in range(max_pages)
    ]
    return [index + 1 for index in dict.fromkeys(indices)]


def _pdf_page_count(path: Path) -> int:
    return len(PdfReader(str(path)).pages)


def _sha1(path: Path) -> str:
    digest = hashlib.sha1()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def perceptual_hash(path: Path, hash_size: int = 8) -> str:
    """Return a DCT perceptual hash that is stable across ordinary resizes."""

    sample_size = hash_size * 4
    with Image.open(path) as image:
        grayscale = image.convert("L").resize(
            (sample_size, sample_size),
            Image.Resampling.LANCZOS,
        )
        pixels = list(grayscale.tobytes())

    cosine = [
        [
            math.cos((2 * position + 1) * frequency * math.pi / (2 * sample_size))
            for position in range(sample_size)
        ]
        for frequency in range(hash_size)
    ]
    coefficients: list[float] = []
    for vertical_frequency in range(hash_size):
        for horizontal_frequency in range(hash_size):
            coefficient = 0.0
            for y in range(sample_size):
                row_offset = y * sample_size
                vertical_weight = cosine[vertical_frequency][y]
                coefficient += vertical_weight * sum(
                    pixels[row_offset + x] * cosine[horizontal_frequency][x]
                    for x in range(sample_size)
                )
            coefficients.append(coefficient)

    threshold = statistics.median(coefficients[1:]) if len(coefficients) > 1 else 0
    value = 0
    for coefficient in coefficients:
        value = (value << 1) | int(coefficient > threshold)
    hex_width = math.ceil(len(coefficients) / 4)
    return f"phash:{value:0{hex_width}x}"


def perceptual_hash_distance(first: str, second: str) -> int:
    first_value = first.removeprefix("phash:")
    second_value = second.removeprefix("phash:")
    if len(first_value) != len(second_value):
        return max(len(first_value), len(second_value)) * 4
    try:
        return (int(first_value, 16) ^ int(second_value, 16)).bit_count()
    except ValueError:
        return 0 if first == second else max(len(first_value), 1) * 4


def _asset_fingerprint(path: Path) -> str:
    try:
        return perceptual_hash(path)
    except Exception:
        return f"sha1:{_sha1(path)}"


def _validate_normalized_bbox(bbox: list[float]) -> None:
    if len(bbox) != 4:
        raise ValueError("bbox 必须包含 [x, y, width, height]。")
    x, y, width, height = bbox
    if not all(math.isfinite(item) for item in bbox):
        raise ValueError("bbox 不能包含非有限数值。")
    if x < 0 or y < 0 or width <= 0 or height <= 0:
        raise ValueError("bbox 坐标必须非负且宽高必须为正。")
    if x > 1 or y > 1 or x + width > 1 + 1e-9 or y + height > 1 + 1e-9:
        raise ValueError("bbox 必须完整位于归一化页面范围内。")


def _asset_url(
    render_id: str,
    filename: str,
    public_base_url: str,
    asset_token: str,
) -> str:
    path = f"/v1/mindmap/assets/{render_id}/{quote(filename)}"
    if not public_base_url:
        return path
    parts = urlsplit(public_base_url)
    safe_base_url = urlunsplit(
        (
            parts.scheme,
            parts.netloc,
            parts.path.rstrip("/"),
            "",
            "",
        )
    )
    return f"{safe_base_url}{path}"


def _run_command(command: list[str], timeout: int = 300) -> None:
    executable = Path(command[0])
    if executable.suffix.lower() in {".cmd", ".bat"}:
        command = [os.environ.get("COMSPEC", "cmd.exe"), "/c", *command]
    subprocess.run(
        command,
        check=True,
        timeout=timeout,
        capture_output=True,
        text=True,
    )


def _find_command(*names: str) -> str | None:
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    return None


def _render_pdf(
    pdf_path: Path,
    render_dir: Path,
    public_base_url: str,
    asset_token: str,
    render_id: str,
    page_numbers: list[int] | None = None,
    resolution: int = 144,
) -> list[RenderedPage]:
    if page_numbers == []:
        return []
    pdftoppm = _find_command("pdftoppm")
    if not pdftoppm:
        raise RuntimeError("未找到 pdftoppm，无法渲染 PDF 页面。")

    generated: list[tuple[int, Path]] = []
    if page_numbers is None:
        prefix = render_dir / "page"
        _run_command(
            [
                pdftoppm,
                "-png",
                "-r",
                str(resolution),
                str(pdf_path),
                str(prefix),
            ]
        )
        generated = [
            (index, source)
            for index, source in enumerate(
                sorted(render_dir.glob("page-*.png")),
                start=1,
            )
        ]
    else:
        for page_number in page_numbers:
            prefix = render_dir / f"selected_{page_number:04d}"
            _run_command(
                [
                    pdftoppm,
                    "-png",
                    "-r",
                    str(resolution),
                    "-f",
                    str(page_number),
                    "-l",
                    str(page_number),
                    "-singlefile",
                    str(pdf_path),
                    str(prefix),
                ]
            )
            source = prefix.with_suffix(".png")
            if not source.exists():
                raise RuntimeError(
                    f"pdftoppm 没有生成第 {page_number} 页图片。"
                )
            generated.append((page_number, source))

    pages: list[RenderedPage] = []
    for page_number, source in generated:
        filename = f"page_{page_number:04d}.png"
        target = render_dir / filename
        source.replace(target)
        with Image.open(target) as image:
            width, height = image.size
        pages.append(
            RenderedPage(
                asset_id=f"page_{page_number:04d}",
                render_id=render_id,
                filename=filename,
                url=_asset_url(
                    render_id,
                    filename,
                    public_base_url,
                    asset_token,
                ),
                page=page_number,
                width=width,
                height=height,
            )
        )
    if not pages:
        raise RuntimeError("pdftoppm 没有生成页面图片。")
    return pages


def _render_pptx(
    pptx_path: Path,
    render_dir: Path,
    public_base_url: str,
    asset_token: str,
    render_id: str,
    page_numbers: list[int] | None = None,
    resolution: int = 144,
) -> list[RenderedPage]:
    soffice = _find_command("soffice", "libreoffice")
    if not soffice:
        raise RuntimeError("未找到 LibreOffice，无法渲染 PPTX 幻灯片。")

    _run_command(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(render_dir),
            str(pptx_path),
        ]
    )
    pdf_path = render_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        pdf_candidates = list(render_dir.glob("*.pdf"))
        if not pdf_candidates:
            raise RuntimeError("LibreOffice 没有生成 PPTX 对应 PDF。")
        pdf_path = pdf_candidates[0]
    try:
        return _render_pdf(
            pdf_path,
            render_dir,
            public_base_url,
            asset_token,
            render_id,
            page_numbers,
            resolution,
        )
    finally:
        pdf_path.unlink(missing_ok=True)


def _normalized_bbox(shape, slide_width: int, slide_height: int) -> list[float]:
    raw_left = shape.left / slide_width
    raw_top = shape.top / slide_height
    raw_right = (shape.left + shape.width) / slide_width
    raw_bottom = (shape.top + shape.height) / slide_height
    left = min(max(raw_left, 0), 1)
    top = min(max(raw_top, 0), 1)
    right = min(max(raw_right, left), 1)
    bottom = min(max(raw_bottom, top), 1)
    return [
        round(left, 6),
        round(top, 6),
        round(right - left, 6),
        round(bottom - top, 6),
    ]


def _crop_image(
    source: Path,
    bbox: list[float],
    target: Path,
    padding_ratio: float = 0.005,
) -> tuple[int, int]:
    _validate_normalized_bbox(bbox)
    with Image.open(source) as image:
        width, height = image.size
        x, y, box_width, box_height = bbox
        pad_x = width * padding_ratio
        pad_y = height * padding_ratio
        left = max(int(x * width - pad_x), 0)
        top = max(int(y * height - pad_y), 0)
        right = min(int((x + box_width) * width + pad_x), width)
        bottom = min(int((y + box_height) * height + pad_y), height)
        if right - left < 4 or bottom - top < 4:
            raise ValueError("裁剪区域过小。")
        cropped = image.crop((left, top, right, bottom)).convert("RGB")
        cropped.save(target, format="PNG")
        return cropped.size


def _table_text(shape) -> str:
    rows: list[str] = []
    for row in shape.table.rows:
        values = [cell.text.strip() for cell in row.cells]
        if any(values):
            rows.append(" | ".join(values))
    return "\n".join(rows)


def _chart_text(shape) -> str:
    chart = shape.chart
    parts: list[str] = []
    if chart.has_title and chart.chart_title.has_text_frame:
        title = chart.chart_title.text_frame.text.strip()
        if title:
            parts.append(title)
    series_names = [
        str(series.name)
        for series in chart.series
        if getattr(series, "name", None)
    ]
    if series_names:
        parts.append("系列：" + "、".join(series_names))
    return "\n".join(parts)


def _extract_pptx_visuals(
    pptx_path: Path,
    render_dir: Path,
    pages: list[RenderedPage],
    public_base_url: str,
    asset_token: str,
    render_id: str,
) -> tuple[list[VisualAsset], list[str]]:
    presentation = Presentation(str(pptx_path))
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    page_by_number = {
        page.page: render_dir / page.filename
        for page in pages
    }
    assets: list[VisualAsset] = []
    warnings: list[str] = []
    canonical_assets: list[VisualAsset] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            bbox = _normalized_bbox(shape, slide_width, slide_height)
            shape_id = getattr(shape, "shape_id", len(assets) + 1)
            base_name = f"slide_{slide_number:04d}_shape_{shape_id}"
            ocr_text = ""
            visual_kind: str | None = None
            filename = ""
            status: str = "ready"
            width: int | None = None
            height: int | None = None

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                extension = shape.image.ext or "png"
                filename = f"{base_name}.{extension}"
                target = render_dir / filename
                target.write_bytes(shape.image.blob)
                try:
                    with Image.open(target) as image:
                        width, height = image.size
                except Exception:
                    width = height = None
                visual_kind = "picture"
            elif getattr(shape, "has_chart", False):
                visual_kind = "chart"
                ocr_text = _chart_text(shape)
            elif getattr(shape, "has_table", False):
                visual_kind = "table"
                ocr_text = _table_text(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                visual_kind = "group_diagram"
            else:
                continue

            if not filename:
                page_path = page_by_number.get(slide_number)
                if page_path and page_path.exists():
                    filename = f"{base_name}.png"
                    target = render_dir / filename
                    try:
                        width, height = _crop_image(page_path, bbox, target)
                    except (OSError, ValueError) as exc:
                        warnings.append(
                            f"{VISUAL_DEGRADED_NATIVE_EXTRACTION} "
                            f"幻灯片 {slide_number} 的 {visual_kind} "
                            f"裁剪失败：{exc}"
                        )
                        continue
                else:
                    status = "needs_render"
                    warnings.append(
                        f"{VISUAL_DEGRADED_NATIVE_EXTRACTION} "
                        f"幻灯片 {slide_number} 的 {visual_kind} "
                        "缺少渲染页，已保留坐标元数据。"
                    )

            target_path = render_dir / filename if filename else None
            fingerprint = (
                _asset_fingerprint(target_path)
                if target_path and target_path.exists()
                else ""
            )
            canonical = next(
                (
                    previous
                    for previous in canonical_assets
                    if fingerprint
                    and previous.sha1
                    and perceptual_hash_distance(
                        fingerprint,
                        previous.sha1,
                    )
                    <= 3
                ),
                None,
            )
            if canonical:
                if canonical.source_slide == slide_number:
                    warnings.append(
                        f"幻灯片 {slide_number} 的 {visual_kind} 与同页已有视觉资产"
                        "感知重复，已跳过重复副本。"
                    )
                    if target_path:
                        target_path.unlink(missing_ok=True)
                    continue
                warnings.append(
                    f"幻灯片 {slide_number} 的 {visual_kind} 与已有视觉资产感知重复，"
                    f"已复用 {canonical.asset_id} 并保留出现位置。"
                )
                if target_path:
                    target_path.unlink(missing_ok=True)
                filename = canonical.filename
                fingerprint = canonical.sha1
                status = canonical.status
            asset_id = f"native_{slide_number:04d}_{shape_id}"
            asset = VisualAsset(
                asset_id=asset_id,
                render_id=render_id,
                filename=filename,
                url=(
                    canonical.url
                    if canonical
                    else _asset_url(
                        render_id,
                        filename,
                        public_base_url,
                        asset_token,
                    )
                    if filename
                    else ""
                ),
                source_slide=slide_number,
                bbox=bbox,
                width=width,
                height=height,
                visual_kind=visual_kind,
                status=status,
                ocr_text=ocr_text,
                # Kept in the legacy field for graph-version compatibility.
                sha1=fingerprint,
            )
            assets.append(asset)
            if fingerprint and canonical is None:
                canonical_assets.append(asset)
    return assets, warnings


def render_document(
    source_path: Path,
    original_filename: str,
    data_root: Path,
    public_base_url: str = "",
    asset_token: str = "",
    max_pages: int | None = None,
    pdf_dpi: int = 144,
) -> RenderResponse:
    suffix = source_path.suffix.lower()
    if suffix not in RENDER_TYPES:
        raise ValueError("视觉渲染仅支持 PDF、PPTX、PNG、JPG、JPEG 或 WEBP。")

    render_id = uuid.uuid4().hex[:16]
    render_dir = data_root / "assets" / render_id
    render_dir.mkdir(parents=True, exist_ok=False)
    warnings: list[str] = []
    pages: list[RenderedPage] = []

    def budgeted_page_numbers(total_pages: int) -> list[int] | None:
        if max_pages is None:
            return None
        selected = _select_page_numbers(total_pages, max(int(max_pages), 0))
        if len(selected) < total_pages:
            labels = "、".join(str(page) for page in selected) or "无"
            warnings.append(
                f"{VISUAL_DEGRADED_RENDER_BUDGET} "
                f"文档共 {total_pages} 页，全页栅格预算为 {max_pages} 页；"
                f"仅栅格化第 {labels} 页，跳过 {total_pages - len(selected)} 页。"
            )
        return selected

    if suffix == ".pdf":
        try:
            page_numbers = (
                budgeted_page_numbers(_pdf_page_count(source_path))
                if max_pages is not None
                else None
            )
            pages = _render_pdf(
                source_path,
                render_dir,
                public_base_url,
                asset_token,
                render_id,
                page_numbers,
                max(int(pdf_dpi), 72),
            )
        except Exception as exc:
            warnings.append(
                f"{VISUAL_DEGRADED_RENDER_FAILURE} "
                f"PDF 页面渲染失败：{exc}"
            )
    elif suffix == ".pptx":
        try:
            page_numbers = (
                budgeted_page_numbers(
                    len(Presentation(str(source_path)).slides)
                )
                if max_pages is not None
                else None
            )
            if page_numbers == []:
                pages = []
            else:
                pages = _render_pptx(
                    source_path,
                    render_dir,
                    public_base_url,
                    asset_token,
                    render_id,
                    page_numbers,
                    max(int(pdf_dpi), 72),
                )
        except Exception as exc:
            warnings.append(
                f"{VISUAL_DEGRADED_RENDER_FAILURE} "
                f"PPTX 整页渲染不可用：{exc}"
            )
    else:
        filename = f"page_0001{suffix}"
        target = render_dir / filename
        shutil.copy2(source_path, target)
        with Image.open(target) as image:
            width, height = image.size
        pages = [
            RenderedPage(
                asset_id="page_0001",
                render_id=render_id,
                filename=filename,
                url=_asset_url(
                    render_id,
                    filename,
                    public_base_url,
                    asset_token,
                ),
                page=1,
                width=width,
                height=height,
            )
        ]

    native_visuals: list[VisualAsset] = []
    if suffix == ".pptx":
        extracted, extraction_warnings = _extract_pptx_visuals(
            source_path,
            render_dir,
            pages,
            public_base_url,
            asset_token,
            render_id,
        )
        native_visuals.extend(extracted)
        warnings.extend(extraction_warnings)

    response = RenderResponse(
        render_id=render_id,
        filename=original_filename,
        pages=pages,
        native_visuals=native_visuals,
        warnings=list(dict.fromkeys(warnings)),
    )
    manifest = {
        **response.model_dump(mode="json"),
        "source_type": suffix.lstrip("."),
        "pdf_dpi": max(int(pdf_dpi), 72),
    }
    (render_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return response


def _crop_context(
    request: CropRequest,
    data_root: Path,
) -> tuple[Path, dict[int, Path], str]:
    render_dir = data_root / "assets" / request.render_id
    manifest_path = render_dir / "manifest.json"
    if not manifest_path.exists():
        raise FileNotFoundError("视觉渲染任务不存在。")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    pages = {
        int(page["page"]): render_dir / page["filename"]
        for page in manifest.get("pages", [])
    }
    source_type = manifest.get("source_type", "")
    return render_dir, pages, source_type



def _crop_visual_unit(
    *,
    render_id: str,
    region: VisualRegion,
    index: int,
    render_dir: Path,
    pages: dict[int, Path],
    source_type: str,
    public_base_url: str,
    asset_token: str,
) -> VisualUnit:
    source = pages.get(region.page)
    if not source or not source.exists():
        raise FileNotFoundError(f"第 {region.page} 页的渲染图片不存在。")
    filename = f"crop_{index:04d}.png"
    target = render_dir / filename
    try:
        width, height = _crop_image(source, region.bbox, target)
    except Exception:
        target.unlink(missing_ok=True)
        raise
    asset_id = f"crop_{index:04d}"
    asset = VisualAsset(
        asset_id=asset_id,
        render_id=render_id,
        filename=filename,
        url=_asset_url(
            render_id,
            filename,
            public_base_url,
            asset_token,
        ),
        source_page=region.page if source_type != "pptx" else None,
        source_slide=region.page if source_type == "pptx" else None,
        bbox=region.bbox,
        width=width,
        height=height,
        visual_kind=region.visual_kind,
        status="ready",
        ocr_text=region.ocr_text,
        # The schema retains the historical name, while the value is now a
        # perceptual fingerprint used by downstream duplicate suppression.
        sha1=_asset_fingerprint(target),
    )
    return VisualUnit(
        asset=asset,
        summary=region.summary,
        knowledge_claims=region.knowledge_claims,
    )


def crop_regions(
    request: CropRequest,
    data_root: Path,
    public_base_url: str = "",
    asset_token: str = "",
) -> list[VisualUnit]:
    render_dir, pages, source_type = _crop_context(request, data_root)
    return [
        _crop_visual_unit(
            render_id=request.render_id,
            region=region,
            index=index,
            render_dir=render_dir,
            pages=pages,
            source_type=source_type,
            public_base_url=public_base_url,
            asset_token=asset_token,
        )
        for index, region in enumerate(request.regions, start=1)
    ]


def crop_regions_best_effort(
    request: CropRequest,
    data_root: Path,
    public_base_url: str = "",
    asset_token: str = "",
) -> tuple[list[VisualUnit | None], list[str]]:
    """Crop independently so one bad model bbox cannot abort the page batch."""

    render_dir, pages, source_type = _crop_context(request, data_root)
    results: list[VisualUnit | None] = []
    warnings: list[str] = []
    for index, region in enumerate(request.regions, start=1):
        try:
            results.append(
                _crop_visual_unit(
                    render_id=request.render_id,
                    region=region,
                    index=index,
                    render_dir=render_dir,
                    pages=pages,
                    source_type=source_type,
                    public_base_url=public_base_url,
                    asset_token=asset_token,
                )
            )
        except (FileNotFoundError, OSError, ValueError) as exc:
            results.append(None)
            warnings.append(
                f"第 {region.page} 页区域 {region.bbox} 裁剪失败：{exc}"
            )
    return results, warnings


def resolve_asset_path(data_root: Path, render_id: str, filename: str) -> Path:
    if not render_id.isalnum():
        raise FileNotFoundError("非法视觉任务 ID。")
    if Path(filename).name != filename:
        raise FileNotFoundError("非法资产文件名。")
    render_dir = (data_root / "assets" / render_id).resolve()
    target = (render_dir / filename).resolve()
    if target.parent != render_dir or not target.is_file():
        raise FileNotFoundError("视觉资产不存在。")
    return target
