from __future__ import annotations

import tempfile
import unittest
from datetime import UTC, datetime
from pathlib import Path
from unittest.mock import patch

from fastapi.testclient import TestClient
from pydantic import ValidationError
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from backend.tests.vnext_test_support import (
    accepted_concept,
    accepted_relation,
    artifact_ref,
    digest,
    graph,
)
from backend.vnext.adapters import LegacyAdaptationBlocked, to_legacy_result
from backend.vnext.api import ShadowAPISettings, create_shadow_app
from backend.vnext.artifacts.canonical import payload_digest
from backend.vnext.artifacts.local_store import LocalArtifactStore
from backend.vnext.canonical_graph import build_canonical_explicit_graph
from backend.vnext.contracts.common import (
    ArtifactProducerRef,
    ArtifactType,
    RuntimeRole,
    StringValue,
)
from backend.vnext.contracts.control import (
    DeclaredRunManifest,
    EvidenceMode,
    ExecutionStatus,
    PublicationStatus,
    QualityAttestation,
    QualityGateDecision,
    QualityMetric,
    QualityStatus,
    RunBudget,
    RunManifest,
    RunProfile,
    evaluate_quality_gate,
)
from backend.vnext.contracts.graph import RelationAssessmentLedger
from backend.vnext.contracts.projection import ProjectionQualityStatus
from backend.vnext.contracts.quality import PilotGateDecision
from backend.vnext.contracts.regions import (
    ReplanAction,
    ReplanRequest,
    ReplanStatus,
)
from backend.vnext.contracts.release import (
    CanaryObservation,
    CanaryStage,
    ReleaseReadinessEvidence,
)
from backend.vnext.orchestration.control_store import SQLiteControlStore
from backend.vnext.orchestration.durable_pipeline import (
    run_durable_shadow_pipeline,
)
from backend.vnext.orchestration.release import (
    ReleaseGateBlocked,
    ReleaseGovernor,
)
from backend.vnext.orchestration.shadow_pipeline import run_shadow_pipeline
from backend.vnext.orchestration.source_shadow import run_source_shadow
from backend.vnext.projection import build_diagnostic_projection
from backend.vnext.source_ir import parse_source


NOW = datetime(2026, 7, 30, 5, 0, tzinfo=UTC)
VALID_SOURCE = (
    "# Carbonyl Chemistry\n"
    "## Foundations\n"
    "Aldehydes are terminal carbonyl compounds.\n"
    "## Applications\n"
    "Ketones are used in synthesis.\n"
)


def _manifest(candidate_status: PublicationStatus) -> RunManifest:
    return RunManifest(
        manifest_id=f"run_manifest_{'1' * 32}",
        run_id=f"run_{'2' * 32}",
        revision=1,
        owner_id="owner-a",
        declared=DeclaredRunManifest(
            source_hash=digest("1"),
            profile=RunProfile.STANDARD,
            evidence_mode=EvidenceMode.SOURCE_ONLY,
            no_egress=True,
            budget=RunBudget(
                max_wall_seconds=600,
                max_model_calls=0,
                max_search_queries=0,
                max_search_fetches=0,
                max_cost_microunits=0,
                vlm_concurrency=0,
                text_concurrency=1,
                search_concurrency=0,
            ),
            code_revision="q0-red-team",
            dependency_digest=digest("2"),
            parser_policy_digest=digest("3"),
            renderer_policy_digest=digest("4"),
            prompt_policy_digest=digest("5"),
            tool_policy_digest=digest("6"),
            search_policy_digest=digest("7"),
            schema_digests=(StringValue(key="test", value=digest("8")),),
            random_seed=0,
        ),
        execution_status=ExecutionStatus.SUCCEEDED,
        quality_status=QualityStatus.PASSED,
        publication_status=candidate_status,
        created_at=NOW,
        updated_at=NOW,
    )


def _release_evidence() -> ReleaseReadinessEvidence:
    return ReleaseReadinessEvidence(
        release_id=f"release_{'3' * 32}",
        owner_id="owner-a",
        run_id=f"run_{'2' * 32}",
        candidate_projection_ref=artifact_ref(
            ArtifactType.DIAGNOSTIC_PROJECTION,
            "4",
        ),
        pilot_report_digest=digest("9"),
        pilot_gate_decision=PilotGateDecision.PASS,
        public_api_approved=True,
        diagnostic_ux_passed=True,
        rollback_drill_passed=True,
        blind_set_expanded=True,
        disaster_recovery_passed=True,
        created_at=NOW,
    )


def _canary_observation() -> CanaryObservation:
    return CanaryObservation(
        release_id=f"release_{'3' * 32}",
        stage=CanaryStage.SHADOW,
        cumulative_samples=0,
        severe_errors=0,
        gate_bypasses=0,
        cross_owner_reads=0,
        rollback_failures=0,
        quality_failed_public_results=0,
        observed_at=NOW,
    )


def _forced_replan(planning, inventory, _audit):
    entry = next(
        item
        for item in inventory.all_entries()
        if item.evidence_refs
    )
    return (
        ReplanRequest(
            request_id=f"replan_{'5' * 32}",
            affected_region_id=planning.root_region_id,
            minimum_replan_ancestor_id=planning.root_region_id,
            omitted_source_ids=(entry.source_id,),
            requested_action=ReplanAction.RESPLIT,
            evidence_refs=entry.evidence_refs,
        ),
    )


class VNextQ0IntegratedRedTeamTests(unittest.TestCase):
    def test_open_replan_blocks_shadow_and_durable_teaching_results(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "course.md"
            source_path.write_text(VALID_SOURCE, encoding="utf-8")

            with patch(
                "backend.vnext.orchestration.shadow_pipeline."
                "audit_regions_bottom_up",
                side_effect=_forced_replan,
            ):
                shadow = run_shadow_pipeline(
                    source_path,
                    owner_id="owner-a",
                    store=LocalArtifactStore(root / "shadow-artifacts"),
                )
            self.assertEqual(
                shadow.projection.quality_status,
                ProjectionQualityStatus.BLOCKED_SEMANTIC,
            )
            self.assertTrue(
                any(
                    "open_replan_quarantined" in item.reason_codes
                    for item in shadow.canonical_graph.unresolved_items
                )
            )

            with patch(
                "backend.vnext.orchestration.durable_pipeline."
                "audit_regions_bottom_up",
                side_effect=_forced_replan,
            ):
                durable = run_durable_shadow_pipeline(
                    source_path,
                    owner_id="owner-a",
                    artifact_store=LocalArtifactStore(
                        root / "durable-artifacts"
                    ),
                    control_store=SQLiteControlStore(
                        root / "control.sqlite3"
                    ),
                    worker_id="q0-red-team",
                )
            self.assertEqual(
                durable.run_manifest.quality_status,
                QualityStatus.BLOCKED_SEMANTIC,
            )
            self.assertEqual(
                durable.run_manifest.publication_status,
                PublicationStatus.DRAFT,
            )
            reused = run_durable_shadow_pipeline(
                source_path,
                owner_id="owner-a",
                artifact_store=LocalArtifactStore(
                    root / "durable-artifacts"
                ),
                control_store=SQLiteControlStore(
                    root / "control.sqlite3"
                ),
                worker_id="q0-red-team-reuse",
                run_id=f"run_{'9' * 32}",
            )
            self.assertEqual(
                reused.run_manifest.quality_status,
                QualityStatus.BLOCKED_SEMANTIC,
            )
            self.assertIn(
                "omission-and-region-audit",
                reused.reused_stages,
            )

    def test_quality_pass_cannot_contain_a_failed_hard_metric(self):
        with self.assertRaisesRegex(
            ValidationError,
            "gate_decision",
        ):
            QualityAttestation(
                attestation_id=f"attestation_{'6' * 32}",
                owner_id="owner-a",
                artifact_ref=artifact_ref(
                    ArtifactType.DIAGNOSTIC_PROJECTION,
                    "7",
                ),
                evaluator=ArtifactProducerRef(
                    producer_id="q0-quality-auditor",
                    producer_version="1.0.0",
                    role=RuntimeRole.QUALITY_AUDITOR,
                ),
                policy_digest=digest("8"),
                closure_digest=digest("9"),
                evaluator_build_digest=digest("a"),
                metrics=(
                    QualityMetric(
                        name="open_replan_count",
                        value=1,
                        threshold=0,
                        passed=False,
                    ),
                ),
                gate_decision=QualityGateDecision.PASS,
                created_at=NOW,
            )

    def test_quality_missing_threshold_is_incomplete_and_store_rechecks(self):
        self.assertEqual(
            evaluate_quality_gate(()),
            QualityGateDecision.INCOMPLETE,
        )
        self.assertEqual(
            evaluate_quality_gate(
                (
                    QualityMetric(
                        name="threshold_not_frozen",
                        value=1,
                    ),
                )
            ),
            QualityGateDecision.INCOMPLETE,
        )
        valid = QualityAttestation(
            attestation_id=f"attestation_{'a' * 32}",
            owner_id="owner-a",
            artifact_ref=artifact_ref(
                ArtifactType.DIAGNOSTIC_PROJECTION,
                "b",
            ),
            evaluator=ArtifactProducerRef(
                producer_id="q0-store-quality-auditor",
                producer_version="1.0.0",
                role=RuntimeRole.QUALITY_AUDITOR,
            ),
            policy_digest=digest("c"),
            closure_digest=digest("d"),
            evaluator_build_digest=digest("e"),
            metrics=(
                QualityMetric(
                    name="hard_failure",
                    value=1,
                    threshold=0,
                    passed=False,
                ),
            ),
            gate_decision=QualityGateDecision.BLOCK,
            created_at=NOW,
        )
        forged = valid.model_copy(
            update={"gate_decision": QualityGateDecision.PASS}
        )
        with tempfile.TemporaryDirectory() as tmp:
            with self.assertRaisesRegex(ValidationError, "gate_decision"):
                SQLiteControlStore(
                    Path(tmp) / "control.sqlite3"
                ).record_quality_attestation(forged)

    def test_raw_manifest_detects_a_page_dropped_by_the_main_parser(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "two-pages.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=612, height=792)
            writer.add_blank_page(width=612, height=792)
            with source_path.open("wb") as handle:
                writer.write(handle)
            parsed = parse_source(source_path)
            dropped = parsed.model_copy(update={"pages": parsed.pages[:1]})

            with patch(
                "backend.vnext.orchestration.source_shadow.parse_source",
                return_value=dropped,
            ):
                result = run_source_shadow(
                    source_path,
                    owner_id="owner-a",
                    store=LocalArtifactStore(root / "artifacts"),
                )

            self.assertEqual(
                result.source_inventory.raw_manifest.native_page_count,
                2,
            )
            self.assertEqual(
                result.source_inventory.raw_manifest.parser_page_count,
                1,
            )
            self.assertIn(
                "native_parser_page_count_mismatch",
                result.source_inventory.raw_manifest.mismatch_codes,
            )
            self.assertTrue(result.source_inventory.unresolved_entries)

    def test_raw_manifest_detects_a_pptx_object_dropped_by_the_parser(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "objects.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[5]
            )
            slide.shapes.title.text = "Object inventory"
            detail = slide.shapes.add_textbox(
                Inches(1.0),
                Inches(1.5),
                Inches(4.0),
                Inches(0.5),
            )
            detail.text = "A source fact"
            presentation.save(source_path)
            parsed = parse_source(source_path)
            page = parsed.pages[0]
            dropped = parsed.model_copy(
                update={
                    "pages": (
                        page.model_copy(
                            update={
                                "native_objects": page.native_objects[:-1],
                            }
                        ),
                    ),
                }
            )

            with patch(
                "backend.vnext.orchestration.source_shadow.parse_source",
                return_value=dropped,
            ):
                result = run_source_shadow(
                    source_path,
                    owner_id="owner-a",
                    store=LocalArtifactStore(root / "artifacts"),
                )

            raw = result.source_inventory.raw_manifest
            self.assertGreater(
                raw.native_object_count or 0,
                raw.parser_object_count,
            )
            self.assertIn(
                "native_parser_object_count_mismatch",
                raw.mismatch_codes,
            )
            self.assertTrue(
                any(
                    item.declared_role
                    == "native_parser_object_count_mismatch"
                    for item in result.source_inventory.unresolved_entries
                )
            )

    def test_raw_manifest_quarantines_unrepresented_pptx_signals(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "signals.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[5]
            )
            slide.shapes.title.text = "Hidden source"
            off_slide = slide.shapes.add_textbox(
                Inches(11.0),
                Inches(1.0),
                Inches(1.0),
                Inches(0.5),
            )
            off_slide.text = "Off-slide evidence"
            off_slide._element.nvSpPr.cNvPr.set(
                "descr",
                "Alternative evidence text",
            )
            slide.notes_slide.notes_text_frame.text = "Teacher note"
            slide._element.set("show", "0")
            presentation.save(source_path)

            result = run_source_shadow(
                source_path,
                owner_id="owner-a",
                store=LocalArtifactStore(root / "artifacts"),
            )

            raw = result.source_inventory.raw_manifest
            self.assertEqual(raw.hidden_page_count, 1)
            self.assertEqual(raw.parser_hidden_page_count, 0)
            self.assertEqual(raw.notes_count, 1)
            self.assertEqual(raw.parser_notes_count, 0)
            self.assertEqual(raw.alt_text_count, 1)
            self.assertEqual(raw.parser_alt_text_count, 0)
            self.assertEqual(raw.off_slide_object_count, 1)
            self.assertEqual(raw.parser_off_slide_object_count, 1)
            self.assertTrue(
                {
                    "native_parser_hidden_page_state_mismatch",
                    "native_parser_notes_count_mismatch",
                    "native_parser_alt_text_count_mismatch",
                }.issubset(raw.mismatch_codes)
            )
            unresolved_roles = {
                item.declared_role
                for item in result.source_inventory.unresolved_entries
            }
            self.assertTrue(set(raw.mismatch_codes).issubset(unresolved_roles))

    def test_closed_replan_requires_new_closure_and_tree_revision(self):
        open_request = ReplanRequest(
            request_id=f"replan_{'c' * 32}",
            affected_region_id=f"reg_{'d' * 32}",
            minimum_replan_ancestor_id=f"reg_{'d' * 32}",
            omitted_source_ids=(
                f"src:block:{'e' * 64}",
            ),
            requested_action=ReplanAction.RESPLIT,
            evidence_refs=(),
        )
        payload = open_request.model_dump(mode="python")
        payload["status"] = ReplanStatus.REJECTED
        with self.assertRaisesRegex(ValidationError, "closure"):
            ReplanRequest.model_validate(payload)
        payload.update(
            {
                "closure_digest": digest("f"),
                "resolved_tree_revision": 2,
            }
        )
        closed = ReplanRequest.model_validate(payload)
        self.assertEqual(closed.status, ReplanStatus.REJECTED)

    def test_owner_is_derived_from_the_authenticated_principal(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            ingest = root / "ingest"
            ingest.mkdir()
            (ingest / "course.md").write_text(
                VALID_SOURCE,
                encoding="utf-8",
            )
            settings = ShadowAPISettings(
                enabled=True,
                service_token="shadow-secret",
                principal_subject="service:test",
                principal_tenant="tenant-a",
                principal_owner_id="owner-a",
                principal_audience="zlb-vnext-shadow",
                principal_scopes=("vnext:run", "vnext:read"),
                ingest_root=ingest,
                artifact_root=root / "artifacts",
                control_db=root / "control.sqlite3",
            )
            application = create_shadow_app(settings)
            client = TestClient(application)

            created = client.post(
                "/v1/shadow/runs",
                json={"source_path": "course.md"},
                headers={
                    "Authorization": "Bearer shadow-secret",
                    "X-VNext-Owner": "owner-b",
                },
            )

            self.assertEqual(created.status_code, 403)
            self.assertTrue(application.state.security_events)
            self.assertEqual(
                application.state.security_events[-1]["code"],
                "owner_header_mismatch",
            )

    def test_governor_rejects_unstored_caller_readiness(self):
        with tempfile.TemporaryDirectory() as tmp:
            governor = ReleaseGovernor(
                SQLiteControlStore(Path(tmp) / "control.sqlite3")
            )

            with self.assertRaises(ReleaseGateBlocked) as raised:
                governor.activate_candidate(
                    _manifest(PublicationStatus.RELEASE_CANDIDATE),
                    _release_evidence(),
                    _canary_observation(),
                    expected_pointer_version=None,
                    expected_release_sequence=None,
                )

            self.assertIn(
                "trusted_release_evidence_missing",
                raised.exception.decision.reason_codes,
            )

    def test_governor_rejects_observation_not_written_by_aggregator(self):
        with tempfile.TemporaryDirectory() as tmp:
            control = SQLiteControlStore(Path(tmp) / "control.sqlite3")
            governor = ReleaseGovernor(control)
            manifest = _manifest(PublicationStatus.RELEASE_CANDIDATE)
            evidence = _release_evidence()
            observation = _canary_observation()
            control.create_run(manifest)
            control.record_release_readiness_evidence(
                evidence,
                recorder=ArtifactProducerRef(
                    producer_id="q0-release-evidence-aggregator",
                    producer_version="1.0.0",
                    role=RuntimeRole.RELEASE_EVIDENCE_AGGREGATOR,
                ),
            )
            control.record_canary_observation(
                observation,
                owner_id="owner-a",
                recorder=ArtifactProducerRef(
                    producer_id="q0-canary-observation-aggregator",
                    producer_version="1.0.0",
                    role=RuntimeRole.CANARY_OBSERVATION_AGGREGATOR,
                ),
            )
            forged = observation.model_copy(
                update={"cumulative_samples": 999}
            )

            with self.assertRaises(ReleaseGateBlocked) as raised:
                governor.activate_candidate(
                    manifest,
                    evidence,
                    forged,
                    expected_pointer_version=None,
                    expected_release_sequence=None,
                )

            self.assertIn(
                "trusted_canary_observation_missing",
                raised.exception.decision.reason_codes,
            )

    def test_canonical_builder_requires_an_independent_relation_ledger(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "course.md"
            source_path.write_text(VALID_SOURCE, encoding="utf-8")
            store = LocalArtifactStore(root / "artifacts")
            result = run_shadow_pipeline(
                source_path,
                owner_id="owner-a",
                store=store,
            )
            proposal_ref = store.ref(
                result.relation_proposal_envelope
            )
            self.assertEqual(
                result.relation_assessment_envelope.input_refs,
                (proposal_ref,),
            )
            self.assertEqual(
                result.relation_proposal_envelope.producer.role,
                RuntimeRole.RELATION_PROPOSER,
            )
            self.assertEqual(
                result.relation_assessment_envelope.producer.role,
                RuntimeRole.RELATION_VERIFIER_A,
            )

            with self.assertRaisesRegex(
                TypeError,
                "relation_assessment_ledger",
            ):
                build_canonical_explicit_graph(
                    result.claim_ledger,
                    result.planning,
                    source_observation_ref=store.ref(
                        result.source.source_envelope
                    ),
                    claim_ledger_ref=store.ref(
                        result.claim_ledger_envelope
                    ),
                )
            payload = (
                result.relation_assessment_ledger.model_dump(mode="python")
            )
            payload["verifier"]["producer_id"] = payload["proposer"][
                "producer_id"
            ]
            with self.assertRaisesRegex(
                ValidationError,
                "proposer",
            ):
                RelationAssessmentLedger.model_validate(payload)

    def test_projection_parent_selection_is_relation_order_invariant(self):
        concepts = (
            accepted_concept("1", "Parent A"),
            accepted_concept("2", "Parent B"),
            accepted_concept("3", "Shared child"),
        )
        first_edge = accepted_relation("4", "1", "3")
        second_edge = accepted_relation("5", "2", "3")
        root_edge = accepted_relation("6", "1", "2")
        first_graph = graph(
            concepts,
            (root_edge, first_edge, second_edge),
        )
        second_graph = graph(
            concepts,
            (second_edge, first_edge, root_edge),
        )

        first = build_diagnostic_projection(
            first_graph,
            canonical_graph_ref=artifact_ref(
                ArtifactType.CANONICAL_EXPLICIT_GRAPH,
                "a",
            ).model_copy(
                update={"payload_digest": payload_digest(first_graph)}
            ),
        )
        second = build_diagnostic_projection(
            second_graph,
            canonical_graph_ref=artifact_ref(
                ArtifactType.CANONICAL_EXPLICIT_GRAPH,
                "b",
            ).model_copy(
                update={"payload_digest": payload_digest(second_graph)}
            ),
        )

        self.assertEqual(
            next(
                item.selected_parent_edge_id
                for item in first.parent_selections
                if item.child_concept_id.endswith("3" * 32)
            ),
            next(
                item.selected_parent_edge_id
                for item in second.parent_selections
                if item.child_concept_id.endswith("3" * 32)
            ),
        )
        self.assertEqual(
            first.quality_status,
            ProjectionQualityStatus.REVIEW_REQUIRED,
        )
        self.assertEqual(
            next(
                item.selection_reason
                for item in first.parent_selections
                if item.child_concept_id.endswith("3" * 32)
            ),
            "stable_id_tiebreak_requires_review",
        )

    def test_legacy_adapter_requires_trusted_publication_closure(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "course.md"
            source_path.write_text(VALID_SOURCE, encoding="utf-8")
            result = run_shadow_pipeline(
                source_path,
                owner_id="owner-a",
                store=LocalArtifactStore(root / "artifacts"),
            )

            with self.assertRaisesRegex(
                LegacyAdaptationBlocked,
                "published pointer",
            ):
                to_legacy_result(
                    task_id="legacy-q0-red-team",
                    run_id=f"run_{'8' * 32}",
                    graph_version=1,
                    filename=source_path.name,
                    file_type="md",
                    source=result.source.source_observation,
                    inventory=result.source.source_inventory,
                    claims=result.claim_ledger,
                    omission_audit=result.omission_audit,
                    graph=result.canonical_graph,
                    projection=result.projection,
                    owner_id="owner-a",
                    control_store=SQLiteControlStore(
                        root / "legacy-control.sqlite3"
                    ),
                )


if __name__ == "__main__":
    unittest.main()
