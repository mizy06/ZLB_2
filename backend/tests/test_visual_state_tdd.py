from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from backend.app import visual_analysis
from backend.app.agents import (
    QWEN_LOW_REASONING_TOKEN_RESERVE,
    STRUCTURED_JSON_TIMEOUT_SECONDS,
    RoleRuntime,
    build_branch_plans,
    synthesize_themes,
)
from backend.app.architecture_schemas import ContentUnit
from backend.app.model_provider import ModelProviderError
from backend.app.mindmap_engine import visuals
from backend.app.mindmap_engine.schemas import (
    RenderResponse,
    RenderedPage,
    VisualAsset,
)
from backend.app.schemas import ParsedDocument


class _FakeVisionClient:
    supports_multimodal = True

    def __init__(self, regions: list[dict]):
        self.regions = regions
        self.calls: list[dict] = []

    async def complete_multimodal_json(self, **kwargs):
        self.calls.append(kwargs)
        return {"regions": self.regions}


class _CapturingThemeClient:
    def __init__(self):
        self.prompt = ""

    async def complete_json(self, **kwargs):
        self.prompt = kwargs["user_prompt"]
        return {
            "root_candidates": [
                {
                    "temp_id": "root",
                    "name": "量子力学课程体系",
                    "support_unit_ids": ["active", "deferred"],
                }
            ],
            "branch_topics": [
                {
                    "temp_id": "branch",
                    "name": "量子态演化规律",
                    "support_unit_ids": ["active", "deferred"],
                }
            ],
        }


class _PartialVisionClient:
    supports_multimodal = True

    async def complete_multimodal_json(self, **kwargs):
        prompt = kwargs["user_prompt"]
        if "页码或幻灯片号：2。" in prompt:
            raise ModelProviderError("page two unavailable")
        return {"regions": []}


def _rendered_fixture(
    root: Path,
    *,
    page_count: int = 1,
) -> tuple[RenderResponse, Path]:
    render_id = "visualstatefixture"
    data_root = root / "data"
    render_dir = data_root / "assets" / render_id
    render_dir.mkdir(parents=True)
    pages: list[RenderedPage] = []
    for page_number in range(1, page_count + 1):
        filename = f"page_{page_number:04d}.png"
        target = render_dir / filename
        image = Image.new("RGB", (240, 160), "white")
        draw = ImageDraw.Draw(image)
        draw.rectangle((28, 24, 116, 134), fill="black")
        draw.ellipse((142, 38, 214, 110), fill="gray")
        image.save(target)
        pages.append(
            RenderedPage(
                asset_id=f"page_{page_number:04d}",
                render_id=render_id,
                filename=filename,
                url=f"/{filename}",
                page=page_number,
                width=image.width,
                height=image.height,
            )
        )
    rendered = RenderResponse(
        render_id=render_id,
        filename="fixture.pdf",
        pages=pages,
        native_visuals=[],
    )
    (render_dir / "manifest.json").write_text(
        json.dumps(
            {
                **rendered.model_dump(mode="json"),
                "source_type": "pdf",
            }
        ),
        encoding="utf-8",
    )
    return rendered, data_root


def _runtime(regions: list[dict]) -> RoleRuntime:
    return RoleRuntime(
        provider="fake",
        model="fake-vision",
        client=_FakeVisionClient(regions),
        available=True,
    )


class VisualStateMachineTDDTests(unittest.IsolatedAsyncioTestCase):
    async def test_visual_analysis_uses_bounded_qwen_completion_policy(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(Path(temp_dir))
            client = _FakeVisionClient([])
            runtime = RoleRuntime(
                provider="qwen",
                model="qwen3.8-max-preview",
                client=client,
                available=True,
            )

            await visual_analysis.analyze_visual_pages(
                document_id="doc_visual_policy",
                rendered=rendered,
                text_units=[],
                runtime=runtime,
                data_root=data_root,
                max_pages=1,
            )

        self.assertEqual(len(client.calls), 1)
        call = client.calls[0]
        answer_budget = call["max_tokens"]
        self.assertEqual(call["max_attempts"], 1)
        self.assertEqual(
            call["thinking_budget"],
            QWEN_LOW_REASONING_TOKEN_RESERVE,
        )
        self.assertNotIn("reasoning_effort", call)
        self.assertEqual(
            call["max_completion_tokens"],
            answer_budget + QWEN_LOW_REASONING_TOKEN_RESERVE,
        )
        self.assertEqual(
            call["timeout_seconds"],
            STRUCTURED_JSON_TIMEOUT_SECONDS,
        )

    async def test_deferred_visual_units_never_enter_theme_or_branch_payloads(self):
        client = _CapturingThemeClient()
        runtime = RoleRuntime(
            provider="fake",
            model="fake-theme",
            client=client,
            available=True,
        )
        units = [
            ContentUnit(
                id="active",
                document_id="doc",
                kind="text",
                importance=0.8,
                status="uncovered",
                text="量子态随时间演化",
                evidence_excerpt="量子态随时间演化",
            ),
            ContentUnit(
                id="deferred",
                document_id="doc",
                kind="visual",
                importance=0,
                status="deferred",
                summary="未能裁剪的图形",
                visual_action="standalone_node",
                knowledge_score=0,
            ),
        ]
        document = ParsedDocument(
            document_id="doc",
            filename="course.pdf",
            file_type="pdf",
            title="量子力学课程体系",
            blocks=[],
        )

        plan, used_model, _warnings = await synthesize_themes(
            document,
            units,
            runtime,
        )
        payload = json.loads(client.prompt)
        plans = build_branch_plans(plan, units)

        self.assertTrue(used_model)
        self.assertEqual(
            [item["unit_id"] for item in payload["content_units"]],
            ["active"],
        )
        self.assertEqual(
            plan.root_candidates[0].support_unit_ids,
            ["active"],
        )
        self.assertEqual(
            plan.branch_topics[0].support_unit_ids,
            ["active"],
        )
        self.assertEqual(
            [unit_id for branch in plans for unit_id in branch.unit_ids],
            ["active"],
        )

    async def test_ignore_decoration_is_persisted_as_rejected_content_unit(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(Path(temp_dir))

            assets, units, used_model, _warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual_state",
                    rendered=rendered,
                    text_units=[],
                    runtime=_runtime(
                        [
                            {
                                "page": 1,
                                "bbox": [0.05, 0.05, 0.2, 0.2],
                                "visual_kind": "logo",
                                "action": "ignore_decoration",
                                "summary": "课程模板角标",
                                "knowledge_claims": ["不承载课程知识"],
                            }
                        ]
                    ),
                    data_root=data_root,
                    max_pages=1,
                )
            )

        self.assertTrue(used_model)
        self.assertEqual(assets, [])
        self.assertEqual(len(units), 1)
        unit = units[0]
        self.assertEqual(unit.status, "rejected")
        self.assertEqual(unit.visual_action, "ignore_decoration")
        self.assertEqual(unit.page, 1)
        self.assertEqual(unit.bbox, [0.05, 0.05, 0.2, 0.2])
        self.assertEqual(unit.knowledge_claims, ["不承载课程知识"])
        self.assertEqual(unit.parent_asset_id, rendered.pages[0].asset_id)
        self.assertIsNone(unit.asset_id)
        self.assertGreaterEqual(unit.decorative_score, 0.9)

    async def test_individual_crop_failure_keeps_deferred_unparsed_unit(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(Path(temp_dir))

            assets, units, used_model, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual_state",
                    rendered=rendered,
                    text_units=[],
                    runtime=_runtime(
                        [
                            {
                                "page": 1,
                                "bbox": [0, 0, 0.001, 0.001],
                                "visual_kind": "diagram",
                                "action": "standalone_node",
                                "ocr_text": "A 到 B",
                                "summary": "流程图",
                                "knowledge_claims": ["A 导致 B"],
                            }
                        ]
                    ),
                    data_root=data_root,
                    max_pages=1,
                )
            )

        self.assertTrue(used_model)
        self.assertEqual(assets, [])
        self.assertEqual(len(units), 1)
        unit = units[0]
        self.assertEqual(unit.status, "deferred")
        self.assertEqual(unit.visual_action, "standalone_node")
        self.assertEqual(unit.page, 1)
        self.assertEqual(unit.bbox, [0, 0, 0.001, 0.001])
        self.assertEqual(unit.ocr_text, "A 到 B")
        self.assertEqual(unit.summary, "流程图")
        self.assertEqual(unit.knowledge_claims, ["A 导致 B"])
        self.assertEqual(unit.parent_asset_id, rendered.pages[0].asset_id)
        self.assertIsNone(unit.asset_id)
        self.assertTrue(any("裁剪失败" in warning for warning in warnings))

    async def test_batch_crop_failure_defers_every_decision(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(Path(temp_dir))
            decisions = [
                {
                    "page": 1,
                    "bbox": [0.05, 0.05, 0.35, 0.35],
                    "visual_kind": "diagram",
                    "action": "standalone_node",
                    "summary": "第一块",
                    "knowledge_claims": ["第一条"],
                },
                {
                    "page": 1,
                    "bbox": [0.55, 0.45, 0.35, 0.35],
                    "visual_kind": "chart",
                    "action": "attach_as_media",
                    "summary": "第二块",
                    "knowledge_claims": ["第二条"],
                },
            ]

            with patch.object(
                visual_analysis,
                "crop_regions_best_effort",
                side_effect=ValueError("manifest corrupted"),
            ):
                assets, units, used_model, warnings = (
                    await visual_analysis.analyze_visual_pages(
                        document_id="doc_visual_state",
                        rendered=rendered,
                        text_units=[],
                        runtime=_runtime(decisions),
                        data_root=data_root,
                        max_pages=1,
                    )
                )

        self.assertTrue(used_model)
        self.assertEqual(assets, [])
        self.assertEqual(len(units), 2)
        self.assertTrue(all(unit.status == "deferred" for unit in units))
        self.assertEqual(
            [unit.knowledge_claims for unit in units],
            [["第一条"], ["第二条"]],
        )
        self.assertTrue(any("视觉裁剪批次失败" in warning for warning in warnings))

    async def test_decompose_crop_points_back_to_full_page_parent_asset(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(Path(temp_dir))

            assets, units, used_model, _warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual_state",
                    rendered=rendered,
                    text_units=[],
                    runtime=_runtime(
                        [
                            {
                                "page": 1,
                                "bbox": [0.05, 0.05, 0.45, 0.85],
                                "visual_kind": "compound_diagram",
                                "action": "decompose",
                                "summary": "复合结构中的左侧子图",
                            }
                        ]
                    ),
                    data_root=data_root,
                    max_pages=1,
                )
            )

        self.assertTrue(used_model)
        self.assertEqual(len(assets), 1)
        self.assertEqual(len(units), 1)
        self.assertEqual(units[0].asset_id, assets[0].asset_id)
        self.assertEqual(units[0].parent_asset_id, rendered.pages[0].asset_id)

    async def test_native_and_vlm_crop_are_reconciled_by_perceptual_hash(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(Path(temp_dir))
            render_dir = data_root / "assets" / rendered.render_id
            bbox = [0.05, 0.05, 0.45, 0.85]
            expected_crop = render_dir / "expected-native.png"
            visuals._crop_image(
                render_dir / rendered.pages[0].filename,
                bbox,
                expected_crop,
            )
            fingerprint = visuals.perceptual_hash(expected_crop)
            expected_crop.unlink()
            native = VisualAsset(
                asset_id="native_0001_7",
                render_id=rendered.render_id,
                filename="native-source.png",
                url="/native-source.png",
                source_page=1,
                bbox=bbox,
                width=100,
                height=120,
                visual_kind="group_diagram",
                status="ready",
                sha1=fingerprint,
            )
            rendered = rendered.model_copy(
                update={"native_visuals": [native]}
            )

            assets, units, used_model, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual_state",
                    rendered=rendered,
                    text_units=[],
                    runtime=_runtime(
                        [
                            {
                                "page": 1,
                                "bbox": bbox,
                                "visual_kind": "group_diagram",
                                "action": "standalone_node",
                                "summary": "与原生图形相同",
                                "knowledge_claims": ["原生图表达课程结构"],
                            }
                        ]
                    ),
                    data_root=data_root,
                    max_pages=1,
                )
            )

        self.assertTrue(used_model)
        self.assertEqual(assets, [])
        self.assertEqual(len(units), 1)
        self.assertEqual(units[0].id, f"visual:{native.asset_id}")
        self.assertEqual(units[0].status, "uncovered")
        self.assertEqual(units[0].asset_id, native.asset_id)
        self.assertIsNone(units[0].parent_asset_id)
        self.assertEqual(units[0].perceptual_hash, native.sha1)
        self.assertEqual(units[0].visual_action, "standalone_node")
        self.assertEqual(units[0].summary, "与原生图形相同")
        self.assertEqual(
            units[0].knowledge_claims,
            ["原生图表达课程结构"],
        )
        self.assertGreaterEqual(units[0].knowledge_score, 0.8)
        self.assertTrue(any("原生视觉资产" in warning for warning in warnings))

    async def test_partial_page_analysis_emits_degraded_signal(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(
                Path(temp_dir),
                page_count=2,
            )
            runtime = RoleRuntime(
                provider="fake",
                model="fake-vision",
                client=_PartialVisionClient(),
                available=True,
            )

            assets, units, complete, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual_state",
                    rendered=rendered,
                    text_units=[],
                    runtime=runtime,
                    data_root=data_root,
                    max_pages=2,
                )
            )

        self.assertEqual(assets, [])
        self.assertEqual(units, [])
        self.assertFalse(complete)
        self.assertTrue(
            any(
                "visual_degraded:partial_page_analysis" in warning
                for warning in warnings
            )
        )
        self.assertTrue(any("第 2 页视觉分析失败" in warning for warning in warnings))

    async def test_page_budget_skip_emits_degraded_signal(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(
                Path(temp_dir),
                page_count=3,
            )

            assets, units, complete, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual_state",
                    rendered=rendered,
                    text_units=[],
                    runtime=_runtime([]),
                    data_root=data_root,
                    max_pages=1,
                )
            )

        self.assertEqual(assets, [])
        self.assertEqual(units, [])
        self.assertFalse(complete)
        self.assertTrue(
            any(
                "visual_degraded:page_budget" in warning
                for warning in warnings
            )
        )
        self.assertTrue(any("未分析 2 页" in warning for warning in warnings))

    def test_pdf_render_budget_rasterizes_only_stratified_pages(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "five-pages.pdf"
            writer = PdfWriter()
            for _ in range(5):
                writer.add_blank_page(width=240, height=160)
            with source.open("wb") as handle:
                writer.write(handle)

            commands: list[list[str]] = []

            def fake_run(command: list[str], timeout: int = 300) -> None:
                del timeout
                commands.append(command)
                prefix = Path(command[-1])
                Image.new("RGB", (240, 160), "white").save(
                    prefix.with_suffix(".png")
                )

            with (
                patch.object(
                    visuals,
                    "_find_command",
                    return_value="/fake/pdftoppm",
                ),
                patch.object(visuals, "_run_command", side_effect=fake_run),
            ):
                rendered = visuals.render_document(
                    source,
                    source.name,
                    root / "data",
                    max_pages=2,
                )

        self.assertEqual([page.page for page in rendered.pages], [1, 5])
        self.assertEqual(len(commands), 2)
        self.assertTrue(all("-singlefile" in command for command in commands))
        self.assertEqual(
            [int(command[command.index("-f") + 1]) for command in commands],
            [1, 5],
        )
        self.assertTrue(
            any(
                "visual_degraded:render_budget" in warning
                for warning in rendered.warnings
            )
        )

    async def test_cross_page_crop_duplicate_keeps_occurrence_provenance(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            rendered, data_root = _rendered_fixture(
                Path(temp_dir),
                page_count=2,
            )

            assets, units, complete, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual_state",
                    rendered=rendered,
                    text_units=[],
                    runtime=_runtime(
                        [
                            {
                                "page": 1,
                                "bbox": [0.05, 0.05, 0.45, 0.85],
                                "visual_kind": "diagram",
                                "action": "standalone_node",
                                "summary": "跨页重复示意图",
                            }
                        ]
                    ),
                    data_root=data_root,
                    max_pages=2,
                )
            )

        self.assertTrue(complete)
        self.assertEqual(len(assets), 1)
        self.assertEqual(len(units), 2)
        self.assertEqual(units[0].status, "uncovered")
        self.assertEqual(units[0].page, 1)
        self.assertEqual(units[1].status, "rejected")
        self.assertEqual(units[1].page, 2)
        self.assertEqual(units[1].asset_id, assets[0].asset_id)
        self.assertEqual(units[1].parent_asset_id, assets[0].asset_id)
        self.assertEqual(units[1].perceptual_hash, assets[0].sha1)
        self.assertTrue(any("保留出现位置" in warning for warning in warnings))

    def test_native_cross_slide_duplicate_reuses_asset_and_keeps_occurrence(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image_path = root / "shared.png"
            Image.new("RGB", (120, 80), "navy").save(image_path)

            presentation = Presentation()
            blank = presentation.slide_layouts[6]
            for _ in range(2):
                slide = presentation.slides.add_slide(blank)
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(1),
                    Inches(1),
                    width=Inches(3),
                )
            source = root / "duplicate.pptx"
            presentation.save(source)
            render_dir = root / "render"
            render_dir.mkdir()

            assets, warnings = visuals._extract_pptx_visuals(
                source,
                render_dir,
                [],
                "",
                "",
                "native-duplicate",
            )

        self.assertEqual(len(assets), 2)
        self.assertEqual([asset.source_slide for asset in assets], [1, 2])
        self.assertNotEqual(assets[0].asset_id, assets[1].asset_id)
        self.assertEqual(assets[0].sha1, assets[1].sha1)
        self.assertEqual(assets[0].filename, assets[1].filename)
        self.assertTrue(any("保留出现位置" in warning for warning in warnings))


if __name__ == "__main__":
    unittest.main()
