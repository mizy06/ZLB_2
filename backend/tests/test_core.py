from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pypdf import PdfWriter

from backend.app.chunking import chunk_document
from backend.app.document_parser import (
    parse_document,
    parse_visual_document,
)
from backend.app.graph_builder import build_graph
from backend.app.heuristics import heuristic_extract


SAMPLE = """# 数据库基础

数据库是用于组织和存储数据的系统。

关系数据库包括表、行和列。索引用于提高查询速度。
"""


class CorePipelineTests(unittest.TestCase):
    def test_text_document_reaches_graph(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "course.md"
            path.write_text(SAMPLE, encoding="utf-8")

            document = parse_document(path)
            chunks = chunk_document(document, max_chars=120)
            extractions = [heuristic_extract(chunk) for chunk in chunks]
            nodes, edges, quality = build_graph(extractions)

        self.assertEqual(document.title, "数据库基础")
        self.assertGreaterEqual(len(chunks), 1)
        self.assertGreaterEqual(len(nodes), 2)
        self.assertEqual(quality.node_count, len(nodes))
        self.assertEqual(quality.edge_count, len(edges))
        self.assertTrue(all(node.evidence for node in nodes))

    def test_chunk_ids_are_stable(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "course.md"
            path.write_text(SAMPLE, encoding="utf-8")
            document = parse_document(path)

            first = chunk_document(document)
            second = chunk_document(document)

        self.assertEqual(
            [chunk.id for chunk in first],
            [chunk.id for chunk in second],
        )

    def test_pptx_preserves_slide_location(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "course.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "线性回归"
            slide.placeholders[1].text = "线性回归用于预测连续数值。"
            presentation.save(path)

            document = parse_document(path)
            chunks = chunk_document(document)

        self.assertEqual(document.title, "线性回归")
        self.assertEqual(chunks[0].slide_start, 1)
        self.assertIn("预测连续数值", chunks[0].text)

    def test_visual_pdf_shell_does_not_extract_text_blocks(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "视觉课程.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=320, height=240)
            with path.open("wb") as handle:
                writer.write(handle)

            document = parse_visual_document(path)

        self.assertEqual(document.title, "视觉课程")
        self.assertEqual(document.blocks, [])
        self.assertEqual(document.parse_metadata["pdf_page_count"], 1)
        self.assertEqual(
            document.parse_metadata["pdf_input_mode"],
            "direct_visual_only",
        )
        self.assertFalse(
            document.parse_metadata["pdf_text_extraction_performed"]
        )


if __name__ == "__main__":
    unittest.main()
