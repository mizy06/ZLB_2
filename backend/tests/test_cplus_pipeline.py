from __future__ import annotations

from dataclasses import replace
from io import BytesIO
import sqlite3
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from backend.app.architecture_schemas import (
    ContentUnit,
    ModelSelection,
    ReviewItemView,
    ReviewResolutionRequest,
)
from backend.app.agents import RoleRuntime
from backend.app.blackboard import SQLiteBlackboard
from backend.app.cplus_pipeline import (
    _mark_discarded_content_units,
    _merge_content_units,
    run_cplus_pipeline,
)
from backend.app.export_service import render_mindmap_png
from backend.app.review_service import resolve_review_item
from backend.app.config import settings
from backend.app.mindmap_engine.schemas import RenderResponse
from backend.app.mindmap_engine.schemas import EvidenceRef, NodeCandidateIn
from backend.app.pdf_page_knowledge import PdfPageKnowledgeResult
from backend.app.schemas import ParsedDocument


SAMPLE = """# 机器学习基础

机器学习是让计算机从数据中学习规律的方法。

## 监督学习

监督学习使用带标签样本学习输入到输出的映射。
分类用于预测离散标签，回归用于预测连续数值。

## 无监督学习

无监督学习处理没有标签的数据。
聚类用于发现样本中的群组结构。

## 模型评估

训练集用于拟合模型，验证集用于选择参数，测试集用于评估泛化能力。
"""


async def noop_progress(stage: str, progress: int, message: str) -> None:
    return None


class CPlusPipelineTests(unittest.IsolatedAsyncioTestCase):
    def test_discarded_planning_units_map_back_to_source_units(self):
        units = [
            ContentUnit(
                id="source-a",
                document_id="doc",
                kind="text",
                status="uncovered",
                text="课程主线",
            ),
            ContentUnit(
                id="source-b",
                document_id="doc",
                kind="text",
                status="uncovered",
                text="科普介绍",
            ),
        ]

        updated = _mark_discarded_content_units(
            units,
            {"source-a"},
            {"source-b": "source-a"},
        )

        self.assertEqual(
            [unit.status for unit in updated],
            ["rejected", "rejected"],
        )
        self.assertEqual(
            [unit.status for unit in units],
            ["uncovered", "uncovered"],
            msg="标记丢弃不能原地修改原始账本对象",
        )

    def test_vlm_enrichment_replaces_the_weak_native_visual_unit(self):
        native = ContentUnit(
            id="visual:native-1",
            document_id="doc",
            kind="visual",
            status="uncovered",
            asset_id="native-1",
            visual_kind="picture",
            visual_action="attach_as_media",
            summary="",
            knowledge_score=0.25,
        )
        enriched = native.model_copy(
            update={
                "visual_kind": "group_diagram",
                "visual_action": "standalone_node",
                "summary": "课程结构图",
                "knowledge_claims": ["A 包含 B"],
                "knowledge_score": 0.86,
            }
        )

        merged = _merge_content_units([native], [enriched])

        self.assertEqual(len(merged), 1)
        self.assertEqual(merged[0].id, native.id)
        self.assertEqual(merged[0].summary, "课程结构图")
        self.assertEqual(merged[0].visual_action, "standalone_node")
        self.assertEqual(merged[0].knowledge_score, 0.86)

    async def test_visual_render_receives_the_configured_page_budget(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "course.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[6]
            )
            textbox = slide.shapes.add_textbox(
                Inches(1),
                Inches(1),
                Inches(6),
                Inches(1),
            )
            textbox.text = "机器学习课程"
            presentation.save(path)
            blackboard = SQLiteBlackboard(root / "blackboard.sqlite3")
            rendered = RenderResponse(
                render_id="render-budget-test",
                filename=path.name,
                pages=[],
                native_visuals=[],
                warnings=[],
            )

            with patch(
                "backend.app.cplus_pipeline.render_document",
                return_value=rendered,
            ) as render:
                await run_cplus_pipeline(
                    task_id="task_visual_budget",
                    file_path=path,
                    filename=path.name,
                    model="qwen3.8-max-preview",
                    provider="qwen",
                    mode="standard",
                    use_ai=False,
                    progress=noop_progress,
                    blackboard=blackboard,
                )

        self.assertEqual(
            render.call_args.kwargs["max_pages"],
            settings.vision_max_pages,
        )

    async def test_page_knowledge_seeds_nodes_without_retranscription(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "course.pdf"
            path.write_bytes(b"%PDF fixture")
            data_root = root / "data"
            render_dir = data_root / "assets" / "knowledge-render"
            render_dir.mkdir(parents=True)
            image_path = render_dir / "page_0001.png"
            Image.new("RGB", (320, 240), "white").save(image_path)
            rendered = RenderResponse.model_validate(
                {
                    "render_id": "knowledge-render",
                    "filename": path.name,
                    "pages": [
                        {
                            "asset_id": "page_0001",
                            "render_id": "knowledge-render",
                            "filename": image_path.name,
                            "url": "/page_0001.png",
                            "page": 1,
                            "width": 320,
                            "height": 240,
                        }
                    ],
                    "native_visuals": [],
                    "warnings": [],
                }
            )
            parsed = ParsedDocument(
                document_id="doc_page_knowledge",
                filename=path.name,
                file_type="pdf",
                title="激光课程",
                blocks=[],
                parse_metadata={
                    "pdf_page_count": 1,
                    "pdf_input_mode": "direct_visual_only",
                    "pdf_text_extraction_performed": False,
                },
            )
            unit = ContentUnit(
                id="pdfk:p0001:linewidth",
                document_id=parsed.document_id,
                kind="visual",
                branch_hint="激光",
                importance=0.9,
                text="Δν/ν≈10^-6",
                heading_path=["激光"],
                unit_role="formula",
                evidence_excerpt="Δν/ν≈10^-6",
                page=1,
                bbox=[0.1, 0.2, 0.7, 0.2],
                asset_id="page_0001",
                visual_kind="direct_page_knowledge",
                visual_action="standalone_node",
                summary="谱线相对宽度满足 Δν/ν≈10^-6",
                knowledge_claims=["谱线相对宽度满足 Δν/ν≈10^-6"],
            )
            candidate = NodeCandidateIn(
                temp_id=f"direct:{unit.id}",
                name="谱线相对宽度",
                type="formula",
                role="formula",
                definition="谱线相对宽度满足 Δν/ν≈10^-6",
                origin="explicit",
                confidence=0.98,
                optional=True,
                evidence=[
                    EvidenceRef(
                        unit_id=unit.id,
                        excerpt=unit.evidence_excerpt,
                        page=1,
                        bbox=unit.bbox,
                        asset_id=unit.asset_id,
                    )
                ],
                support_unit_ids=[unit.id],
            )
            knowledge_document = parsed.model_copy(
                update={
                    "blocks": [],
                    "parse_metadata": {
                        **parsed.parse_metadata,
                        "pdf_page_knowledge": {
                            "complete": True,
                            "accepted_pages": [1],
                            "failed_pages": [],
                            "node_count": 1,
                        },
                    },
                }
            )
            knowledge = PdfPageKnowledgeResult(
                document=knowledge_document,
                content_units=[unit],
                node_candidates=[candidate],
                complete=True,
                accepted_pages=[1],
                degraded_pages=[1],
                called_pages=[1],
            )
            unavailable = RoleRuntime(
                provider="qwen",
                model="qwen3.8-max-preview",
                client=None,
                available=False,
            )
            selection = ModelSelection(
                generator_provider="qwen",
                verifier_provider="qwen",
                vision_provider="qwen",
            )
            blackboard = SQLiteBlackboard(root / "blackboard.sqlite3")

            with (
                patch(
                    "backend.app.cplus_pipeline.settings",
                    replace(
                        settings,
                        mindmap_data_dir=data_root,
                        pdf_transcription_mode="vision_nodes_strict",
                        pdf_page_extraction_mode="layout_nodes",
                        pdf_page_knowledge_prompt_version=(
                            "page-knowledge-stage-v1"
                        ),
                    ),
                ),
                patch(
                    "backend.app.cplus_pipeline.parse_visual_document",
                    return_value=parsed,
                ),
                patch(
                    "backend.app.cplus_pipeline.parse_document",
                ) as parse_text,
                patch(
                    "backend.app.cplus_pipeline.render_document",
                    return_value=rendered,
                ),
                patch(
                    "backend.app.cplus_pipeline.extract_pdf_page_knowledge",
                    return_value=knowledge,
                    create=True,
                ) as extract,
                patch(
                    "backend.app.cplus_pipeline.analyze_visual_pages",
                ) as analyze_visual,
                patch(
                    "backend.app.cplus_pipeline.build_role_runtimes",
                    return_value=(
                        unavailable,
                        unavailable,
                        unavailable,
                        None,
                        None,
                        selection,
                        [],
                    ),
                ),
            ):
                result = await run_cplus_pipeline(
                    task_id="task_page_knowledge",
                    file_path=path,
                    filename=path.name,
                    model="qwen3.8-max-preview",
                    provider="qwen",
                    mode="standard",
                    use_ai=True,
                    progress=noop_progress,
                    blackboard=blackboard,
                )

        extract.assert_awaited_once()
        self.assertEqual(
            extract.call_args.kwargs["extraction_profile"],
            "direct",
        )
        self.assertEqual(
            extract.call_args.kwargs["prompt_version"],
            "page-knowledge-stage-v1",
        )
        parse_text.assert_not_called()
        analyze_visual.assert_not_awaited()
        self.assertEqual(result.document.blocks, [])
        self.assertEqual(result.chunks, [])
        self.assertIn(
            "谱线相对宽度",
            {node.name for node in result.nodes},
        )
        self.assertNotIn("10^6", str(result.model_dump(mode="json")))
        self.assertIn(
            "pdf_page_knowledge",
            result.degraded_components,
        )

    async def test_heuristic_pipeline_builds_versioned_rooted_tree(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "course.md"
            path.write_text(SAMPLE, encoding="utf-8")
            database_path = root / "blackboard.sqlite3"
            blackboard = SQLiteBlackboard(database_path)

            result = await run_cplus_pipeline(
                task_id="task_demo",
                file_path=path,
                filename=path.name,
                model="qwen3.8-max-preview",
                provider="qwen",
                mode="standard",
                use_ai=False,
                progress=noop_progress,
                blackboard=blackboard,
            )

            self.assertEqual(result.graph_version, 1)
            self.assertTrue(result.quality_report.topology_valid)
            self.assertEqual(result.quality_report.root_count, 1)
            self.assertEqual(len(result.tree_edges), len(result.nodes) - 1)
            self.assertEqual(result.quality_report.evidence_coverage, 1)
            self.assertGreaterEqual(
                result.quality_report.weighted_content_coverage,
                0.78,
            )
            self.assertTrue(all(node.depth >= 0 for node in result.nodes))
            self.assertEqual(
                [1],
                blackboard.list_graph_versions("task_demo"),
            )
            restored = blackboard.load_latest_result("task_demo")
            self.assertIsNotNone(restored)
            self.assertEqual(restored.root_id, result.root_id)
            history = blackboard.list_history()
            self.assertEqual(len(history), 1)
            self.assertEqual(history[0]["task_id"], "task_demo")
            self.assertEqual(history[0]["title"], result.document.title)
            self.assertEqual(history[0]["node_count"], len(result.nodes))
            self.assertEqual(history[0]["graph_version"], 1)
            reopened = SQLiteBlackboard(database_path)
            self.assertEqual(len(reopened.list_history()), 1)
            self.assertEqual(
                reopened.load_latest_result("task_demo").root_id,
                result.root_id,
            )
            png = render_mindmap_png(result)
            self.assertTrue(png.startswith(b"\x89PNG\r\n\x1a\n"))
            with Image.open(BytesIO(png)) as image:
                aspect_ratio = max(image.size) / min(image.size)
            self.assertLessEqual(aspect_ratio, 2)

    async def test_human_review_creates_new_graph_version(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "course.md"
            path.write_text(SAMPLE, encoding="utf-8")
            blackboard = SQLiteBlackboard(root / "blackboard.sqlite3")
            result = await run_cplus_pipeline(
                task_id="task_review",
                file_path=path,
                filename=path.name,
                model="qwen3.8-max-preview",
                provider="qwen",
                mode="standard",
                use_ai=False,
                progress=noop_progress,
                blackboard=blackboard,
            )
            subject = next(node for node in result.nodes if node.id != result.root_id)
            review = ReviewItemView(
                id="review_manual",
                type="abstract_parent",
                risk_score=0.7,
                subject_ids=[subject.id],
                reason="测试人工改名",
                evidence_unit_ids=[
                    item.unit_id or item.chunk_id
                    for item in subject.evidence
                    if item.unit_id or item.chunk_id
                ],
            )
            augmented = result.model_copy(
                update={
                    "graph_version": 0,
                    "review_items": [*result.review_items, review],
                }
            )
            blackboard.save_review_items(result.run_id, [review])
            augmented_version = blackboard.save_graph_version(
                result.run_id,
                augmented,
            )

            updated = resolve_review_item(
                blackboard=blackboard,
                task_id="task_review",
                review_id=review.id,
                request=ReviewResolutionRequest(
                    action="rename",
                    label="人工确认主题",
                    reason="更符合课程术语",
                    expected_graph_version=augmented_version,
                ),
            )

            self.assertEqual(updated.graph_version, 3)
            self.assertEqual(
                next(node for node in updated.nodes if node.id == subject.id).name,
                "人工确认主题",
            )
            self.assertEqual(
                next(item for item in updated.review_items if item.id == review.id).status,
                "resolved",
            )
            self.assertEqual(updated.decision_records[-1].actor, "human")

    async def test_model_preflight_failure_marks_run_failed(self):
        async def fail_preflight(**kwargs):
            raise RuntimeError("model preflight failed")

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "course.md"
            path.write_text(SAMPLE, encoding="utf-8")
            database_path = root / "blackboard.sqlite3"
            blackboard = SQLiteBlackboard(database_path)

            with patch(
                "backend.app.cplus_pipeline.build_role_runtimes",
                fail_preflight,
            ):
                with self.assertRaisesRegex(RuntimeError, "model preflight failed"):
                    await run_cplus_pipeline(
                        task_id="task_preflight_failure",
                        file_path=path,
                        filename=path.name,
                        model="qwen3.8-max-preview",
                        provider="qwen",
                        mode="standard",
                        use_ai=True,
                        progress=noop_progress,
                        blackboard=blackboard,
                    )

            with sqlite3.connect(database_path) as connection:
                status, stage = connection.execute(
                    "SELECT status, stage FROM runs WHERE task_id = ?",
                    ("task_preflight_failure",),
                ).fetchone()
            self.assertEqual((status, stage), ("failed", "failed"))


if __name__ == "__main__":
    unittest.main()
