from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from pydantic import ValidationError

from backend.app import document_parser, visual_analysis
from backend.app.agents import RoleRuntime
from backend.app.architecture_schemas import ContentUnit
from backend.app.chunking import chunk_document
from backend.app.document_parser import parse_document
from backend.app.heuristics import heuristic_extract
from backend.app.mindmap_engine import visuals
from backend.app.mindmap_engine.schemas import (
    CropRequest,
    RenderResponse,
    RenderedPage,
    VisualRegion,
)
from backend.app.schemas import Chunk, ParsedDocument, SourceBlock


class _FakeVisionClient:
    supports_multimodal = True

    def __init__(self, regions: list[dict] | None = None):
        self.regions = regions or []
        self.prompts: list[str] = []

    async def complete_multimodal_json(self, **kwargs):
        self.prompts.append(kwargs["user_prompt"])
        return {"regions": self.regions}


def _rendered_fixture(
    root: Path,
    *,
    page_count: int,
    image_factory=None,
) -> tuple[RenderResponse, Path]:
    render_id = "renderfixture01"
    data_root = root / "data"
    render_dir = data_root / "assets" / render_id
    render_dir.mkdir(parents=True)
    pages: list[RenderedPage] = []
    for page_number in range(1, page_count + 1):
        filename = f"page_{page_number:04d}.png"
        target = render_dir / filename
        image = (
            image_factory(page_number)
            if image_factory
            else Image.new("RGB", (160, 100), "white")
        )
        image.save(target)
        pages.append(
            RenderedPage(
                asset_id=f"page_{page_number:04d}",
                render_id=render_id,
                filename=filename,
                url=f"/{filename}",
                page=page_number,
                width=image.width,
                height=image.height,
            )
        )
    response = RenderResponse(
        render_id=render_id,
        filename="fixture.pdf",
        pages=pages,
        native_visuals=[],
    )
    (render_dir / "manifest.json").write_text(
        json.dumps(
            {
                **response.model_dump(mode="json"),
                "source_type": "pdf",
            }
        ),
        encoding="utf-8",
    )
    return response, data_root


class InputPipelineTDDTests(unittest.TestCase):
    def test_chunks_respect_page_heading_boundaries_and_keep_overlap_out_of_text(self):
        first_page_text = "".join(chr(0x4E00 + index) for index in range(55))
        document = ParsedDocument(
            document_id="doc_boundaries",
            filename="course.pdf",
            file_type="pdf",
            title="课程",
            blocks=[
                SourceBlock(text=first_page_text, page=1, heading="第一章"),
                SourceBlock(text="第二页知识点", page=2, heading="第二章"),
            ],
        )

        chunks = chunk_document(document, max_chars=24, overlap_chars=6)
        first_page = [chunk for chunk in chunks if chunk.page_start == 1]
        second_page = [chunk for chunk in chunks if chunk.page_start == 2]

        self.assertGreater(len(first_page), 1)
        self.assertTrue(
            all(chunk.page_start == chunk.page_end == 1 for chunk in first_page)
        )
        self.assertEqual(len(second_page), 1)
        self.assertEqual(second_page[0].page_end, 2)
        self.assertEqual(second_page[0].heading, "第二章")
        self.assertEqual(second_page[0].context_before, "")
        self.assertTrue(first_page[1].context_before)
        self.assertTrue(all(len(chunk.text) <= 24 for chunk in chunks))
        self.assertTrue(
            all("[上文衔接]" not in chunk.text for chunk in chunks)
        )
        self.assertEqual("".join(chunk.text for chunk in first_page), first_page_text)

        restored = Chunk.model_validate_json(first_page[1].model_dump_json())
        self.assertEqual(restored.context_before, first_page[1].context_before)

    def test_legacy_chunk_without_context_field_still_loads(self):
        restored = Chunk.model_validate(
            {
                "id": "legacy_chunk",
                "index": 0,
                "text": "旧图版本正文",
                "page_start": 1,
                "page_end": 1,
            }
        )

        self.assertEqual(restored.context_before, "")

    def test_heading_change_is_a_hard_boundary_without_page_metadata(self):
        document = ParsedDocument(
            document_id="doc_headings",
            filename="course.md",
            file_type="md",
            title="课程",
            blocks=[
                SourceBlock(text="第一章正文", heading="第一章"),
                SourceBlock(text="第二章正文", heading="第二章"),
            ],
        )

        chunks = chunk_document(document, max_chars=100, overlap_chars=10)

        self.assertEqual([chunk.heading for chunk in chunks], ["第一章", "第二章"])
        self.assertEqual([chunk.text for chunk in chunks], ["第一章正文", "第二章正文"])

    def test_pdf_reflow_repairs_hard_wraps_and_english_hyphenation(self):
        raw = (
            "The electro-\n"
            "magnetic field explains waves.\n"
            "中文换行内容\n"
            "继续说明。\n\n"
            "• 第一项\n"
            "• 第二项"
        )

        text = document_parser._reflow_pdf_text(raw)

        self.assertIn("The electromagnetic field explains waves.", text)
        self.assertIn("中文换行内容继续说明。", text)
        self.assertIn("• 第一项\n• 第二项", text)
        self.assertNotIn("electro-\nmagnetic", text)

    def test_pdf_parser_prefers_layout_extraction_mode(self):
        class _Page:
            def __init__(self):
                self.calls: list[dict] = []

            def extract_text(self, **kwargs):
                self.calls.append(kwargs)
                return "版面文本"

        page = _Page()
        reader = type("_Reader", (), {"pages": [page]})()
        with patch.object(document_parser, "PdfReader", return_value=reader):
            blocks = document_parser._parse_pdf(Path("fixture.pdf"))

        self.assertEqual(page.calls, [{"extraction_mode": "layout"}])
        self.assertEqual(blocks[0].page, 1)

    def test_pptx_text_uses_title_then_visual_reading_order_and_keeps_slides(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "ordered.pptx"
            presentation = Presentation()
            first = presentation.slides.add_slide(presentation.slide_layouts[5])
            first.shapes.title.text = "第一页标题"
            lower = first.shapes.add_textbox(
                Inches(1),
                Inches(4),
                Inches(5),
                Inches(0.5),
            )
            lower.text = "下方内容"
            upper = first.shapes.add_textbox(
                Inches(1),
                Inches(2),
                Inches(5),
                Inches(0.5),
            )
            upper.text = "上方内容"
            second = presentation.slides.add_slide(presentation.slide_layouts[5])
            second.shapes.title.text = "第二页标题"
            second.shapes.add_textbox(
                Inches(1),
                Inches(2),
                Inches(5),
                Inches(0.5),
            ).text = "第二页内容"
            presentation.save(path)

            document = parse_document(path)
            chunks = chunk_document(document, max_chars=500)

        first_text = next(block.text for block in document.blocks if block.slide == 1)
        self.assertLess(first_text.index("第一页标题"), first_text.index("上方内容"))
        self.assertLess(first_text.index("上方内容"), first_text.index("下方内容"))
        self.assertEqual(
            [(chunk.slide_start, chunk.slide_end) for chunk in chunks],
            [(1, 1), (2, 2)],
        )

    def test_heuristic_label_gate_rejects_fragments_without_mutating_them(self):
        chunk = Chunk(
            id="chunk_labels",
            index=0,
            text="\n".join(
                [
                    "电子自旋（角动量）",
                    "的光谱结构（即使",
                    "但是",
                    "原子沉积层不",
                    "[上文衔接]",
                    "与管壁碰",
                    "微观粒子的运动状态是用波函数描",
                    "异常\x07标签",
                ]
            ),
            page_start=7,
            page_end=7,
        )

        extraction = heuristic_extract(chunk)
        names = {node.name for node in extraction.nodes}

        self.assertIn("电子自旋（角动量）", names)
        self.assertNotIn("的光谱结构（即使", names)
        self.assertNotIn("但是", names)
        self.assertNotIn("原子沉积层不", names)
        self.assertNotIn("[上文衔接]", names)
        self.assertNotIn("与管壁碰", names)
        self.assertNotIn("微观粒子的运动状态是用波函数描", names)
        self.assertNotIn("异常\x07标签", names)
        self.assertTrue(
            all(
                evidence.page == 7
                for node in extraction.nodes
                for evidence in node.evidence
            )
        )

    def test_visual_bbox_must_fit_inside_the_normalized_page(self):
        with self.assertRaises(ValidationError):
            visual_analysis.VisualRegionDecision(
                page=1,
                bbox=[0.8, 0.8, 0.3, 0.3],
                action="standalone_node",
            )
        with self.assertRaises(ValidationError):
            visual_analysis.VisualRegionDecision(
                page=1,
                bbox=[0, 0, float("nan"), 0.3],
                action="standalone_node",
            )

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.png"
            Image.new("RGB", (100, 100), "white").save(source)
            rendered = visuals.render_document(source, "source.png", root / "data")
            with self.assertRaises(ValueError):
                visuals.crop_regions(
                    CropRequest(
                        render_id=rendered.render_id,
                        regions=[
                            VisualRegion(
                                page=1,
                                bbox=[0.8, 0.8, 0.3, 0.3],
                                visual_kind="diagram",
                            )
                        ],
                    ),
                    root / "data",
                )

    def test_visual_nms_suppresses_duplicate_overlapping_regions(self):
        first = visual_analysis.VisualRegionDecision(
            page=1,
            bbox=[0.1, 0.1, 0.6, 0.6],
            action="standalone_node",
            summary="完整区域",
        )
        duplicate = visual_analysis.VisualRegionDecision(
            page=1,
            bbox=[0.11, 0.11, 0.59, 0.59],
            action="attach_as_media",
            summary="重复区域",
        )

        kept, suppressed = visual_analysis._suppress_overlapping_decisions(
            [first, duplicate]
        )

        self.assertEqual(len(kept), 1)
        self.assertEqual(suppressed, 1)

    def test_perceptual_hash_is_stable_across_resizing(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            small = root / "small.png"
            large = root / "large.png"
            image = Image.new("RGB", (64, 64), "white")
            draw = ImageDraw.Draw(image)
            draw.rectangle((8, 8, 28, 52), fill="black")
            draw.ellipse((34, 16, 56, 38), fill="black")
            image.save(small)
            image.resize((160, 160)).save(large)

            small_hash = visuals.perceptual_hash(small)
            large_hash = visuals.perceptual_hash(large)

        self.assertLessEqual(
            visuals.perceptual_hash_distance(small_hash, large_hash),
            4,
        )

    def test_native_pptx_visuals_drop_resized_duplicate_images(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            small = root / "small.png"
            large = root / "large.jpg"
            image = Image.new("RGB", (64, 64), "white")
            draw = ImageDraw.Draw(image)
            draw.rectangle((8, 8, 28, 52), fill="black")
            draw.ellipse((34, 16, 56, 38), fill="black")
            image.save(small)
            image.resize((160, 160)).save(large, quality=90)

            pptx_path = root / "duplicates.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(str(small), Inches(1), Inches(1))
            slide.shapes.add_picture(str(large), Inches(5), Inches(1))
            presentation.save(pptx_path)
            render_dir = root / "render"
            render_dir.mkdir()

            assets, warnings = visuals._extract_pptx_visuals(
                pptx_path,
                render_dir,
                [],
                "",
                "",
                "renderfixture02",
            )

        self.assertEqual(len(assets), 1)
        self.assertTrue(assets[0].sha1.startswith("phash:"))
        self.assertTrue(any("感知重复" in warning for warning in warnings))


class VisualAnalysisTDDTests(unittest.IsolatedAsyncioTestCase):
    async def test_visual_page_limit_uses_document_wide_stratified_sampling(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            rendered, data_root = _rendered_fixture(root, page_count=6)
            client = _FakeVisionClient()
            runtime = RoleRuntime(
                provider="fake",
                model="fake-vision",
                client=client,
                available=True,
            )
            text_units = [
                ContentUnit(
                    id="page-6-text",
                    document_id="doc_visual",
                    kind="text",
                    text="末页证据",
                    evidence_excerpt="末页证据",
                    page=6,
                )
            ]

            assets, units, used_model, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual",
                    rendered=rendered,
                    text_units=text_units,
                    runtime=runtime,
                    data_root=data_root,
                    max_pages=3,
                )
            )

        self.assertEqual(assets, [])
        self.assertEqual(units, [])
        self.assertFalse(used_model)
        analyzed_pages = [
            int(prompt.split("：", 1)[1].split("。", 1)[0])
            for prompt in client.prompts
        ]
        self.assertEqual(set(analyzed_pages), {1, 3, 6})
        self.assertNotEqual(set(analyzed_pages), {1, 2, 3})
        self.assertTrue(any("分层分析" in warning for warning in warnings))
        self.assertTrue(any("未分析 3 页" in warning for warning in warnings))
        self.assertTrue(
            any(
                "visual_degraded:page_budget" in warning
                for warning in warnings
            )
        )
        page_six_prompt = next(
            prompt
            for prompt in client.prompts
            if "当前页码或幻灯片号：6。" in prompt
        )
        self.assertIn("末页证据", page_six_prompt)

    async def test_crop_failure_is_recorded_without_aborting_visual_stage(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            rendered, data_root = _rendered_fixture(root, page_count=1)
            client = _FakeVisionClient(
                [
                    {
                        "page": 1,
                        "bbox": [0, 0, 0.001, 0.001],
                        "visual_kind": "diagram",
                        "action": "standalone_node",
                        "summary": "过小区域",
                    }
                ]
            )
            runtime = RoleRuntime(
                provider="fake",
                model="fake-vision",
                client=client,
                available=True,
            )

            assets, units, used_model, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual",
                    rendered=rendered,
                    text_units=[],
                    runtime=runtime,
                    data_root=data_root,
                    max_pages=1,
                )
            )

        self.assertEqual(assets, [])
        self.assertEqual(len(units), 1)
        self.assertEqual(units[0].status, "deferred")
        self.assertEqual(units[0].bbox, [0, 0, 0.001, 0.001])
        self.assertEqual(units[0].summary, "过小区域")
        self.assertEqual(units[0].parent_asset_id, rendered.pages[0].asset_id)
        self.assertTrue(used_model)
        self.assertTrue(any("裁剪失败" in warning for warning in warnings))

    async def test_visual_regions_are_deduplicated_by_perceptual_content(self):
        def repeated_image(_page_number: int) -> Image.Image:
            tile = Image.new("RGB", (100, 100), "white")
            draw = ImageDraw.Draw(tile)
            draw.rectangle((10, 10, 40, 85), fill="black")
            draw.ellipse((55, 20, 85, 50), fill="black")
            image = Image.new("RGB", (200, 100), "white")
            image.paste(tile, (0, 0))
            image.paste(tile, (100, 0))
            return image

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            rendered, data_root = _rendered_fixture(
                root,
                page_count=1,
                image_factory=repeated_image,
            )
            client = _FakeVisionClient(
                [
                    {
                        "page": 1,
                        "bbox": [0, 0, 0.5, 1],
                        "visual_kind": "diagram",
                        "action": "attach_as_media",
                        "summary": "左侧重复图",
                    },
                    {
                        "page": 1,
                        "bbox": [0.5, 0, 0.5, 1],
                        "visual_kind": "diagram",
                        "action": "attach_as_media",
                        "summary": "右侧重复图",
                    },
                ]
            )
            runtime = RoleRuntime(
                provider="fake",
                model="fake-vision",
                client=client,
                available=True,
            )

            assets, units, used_model, warnings = (
                await visual_analysis.analyze_visual_pages(
                    document_id="doc_visual",
                    rendered=rendered,
                    text_units=[],
                    runtime=runtime,
                    data_root=data_root,
                    max_pages=1,
                )
            )

        self.assertTrue(used_model)
        self.assertEqual(len(assets), 1)
        self.assertEqual(len(units), 1)
        self.assertEqual(units[0].status, "uncovered")
        self.assertTrue(units[0].perceptual_hash.startswith("phash:"))
        self.assertTrue(any("感知重复" in warning for warning in warnings))


if __name__ == "__main__":
    unittest.main()
