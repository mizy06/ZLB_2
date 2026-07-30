from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from backend.vnext.artifacts.local_store import LocalArtifactStore
from backend.vnext.contracts.common import (
    ArtifactType,
    InterpretationStatus,
)
from backend.vnext.contracts.inventory import (
    InventoryInspectionStatus,
)
from backend.vnext.contracts.source import NativeObjectKind
from backend.vnext.orchestration.source_shadow import run_source_shadow
from backend.vnext.source_ir import parse_source


class VNextSourceShadowTests(unittest.TestCase):
    def test_markdown_ids_outline_and_inventory_are_deterministic(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "course.md"
            source_path.write_text(
                "# Carbonyl Chemistry\n"
                "Aldehydes contain a terminal carbonyl group.\n"
                "## Reactions\n"
                "E = mc^2\n"
                "A + B -> C\n",
                encoding="utf-8",
            )

            first = parse_source(source_path)
            second = parse_source(source_path)

            self.assertEqual(first, second)
            self.assertEqual(len(first.pages), 1)
            self.assertEqual(
                [item.observed_level for item in first.outline_entries],
                [1, 2],
            )
            self.assertTrue(first.pages[0].reading_order_hypotheses)
            self.assertTrue(
                all(
                    item.interpretation_status
                    is InterpretationStatus.INFERRED
                    for item in first.pages[0].role_hypotheses
                )
            )
            kinds = {
                item.kind for item in first.pages[0].native_objects
            }
            self.assertIn(NativeObjectKind.FORMULA, kinds)
            self.assertIn(NativeObjectKind.CHEMICAL_REACTION, kinds)

            store = LocalArtifactStore(root / "shadow")
            result = run_source_shadow(
                source_path,
                owner_id="tenant-a",
                store=store,
            )
            inventory = result.source_inventory
            self.assertEqual(
                inventory.document_ir_ref.artifact_id,
                result.source_envelope.artifact_id,
            )
            self.assertEqual(len(inventory.page_entries), 1)
            self.assertEqual(
                len(inventory.outline_entries),
                len(first.outline_entries),
            )
            self.assertEqual(len(inventory.formula_region_entries), 1)
            self.assertEqual(len(inventory.reaction_region_entries), 1)
            self.assertTrue(
                all(
                    entry.inspection_status
                    is InventoryInspectionStatus.UNRESOLVED
                    for entry in (
                        *inventory.formula_region_entries,
                        *inventory.reaction_region_entries,
                    )
                )
            )
            source_ids = [
                entry.source_id for entry in inventory.all_entries()
            ]
            self.assertEqual(len(source_ids), len(set(source_ids)))

            source_path.write_text(
                source_path.read_text(encoding="utf-8") + "New fact.\n",
                encoding="utf-8",
            )
            changed = parse_source(source_path)
            self.assertNotEqual(changed.source_hash, first.source_hash)
            self.assertNotEqual(changed.document_id, first.document_id)

    def test_pptx_retains_empty_slides_objects_tables_and_headers(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image_path = root / "diagram.png"
            Image.new("RGB", (32, 24), color=(210, 220, 230)).save(
                image_path
            )
            source_path = root / "course.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[5]
            )
            slide.shapes.title.text = "10.1 Carbonyl Compounds"
            formula = slide.shapes.add_textbox(
                Inches(0.8),
                Inches(1.3),
                Inches(2.5),
                Inches(0.5),
            )
            formula.text = "E = mc^2"
            reaction = slide.shapes.add_textbox(
                Inches(0.8),
                Inches(2.0),
                Inches(2.5),
                Inches(0.5),
            )
            reaction.text = "A + B -> C"
            table_shape = slide.shapes.add_table(
                2,
                2,
                Inches(0.8),
                Inches(2.8),
                Inches(4.0),
                Inches(1.2),
            )
            table_shape.table.cell(0, 0).text = "Name"
            table_shape.table.cell(0, 1).text = "Property"
            table_shape.table.cell(1, 0).text = "Aldehyde"
            table_shape.table.cell(1, 1).text = "Terminal"
            slide.shapes.add_picture(
                str(image_path),
                Inches(5.2),
                Inches(1.3),
                width=Inches(1.5),
            )
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(source_path)

            source = parse_source(source_path)

            self.assertEqual(len(source.pages), 2)
            self.assertEqual(len(source.outline_entries), 1)
            first_page = source.pages[0]
            kinds = [item.kind for item in first_page.native_objects]
            self.assertIn(NativeObjectKind.TABLE, kinds)
            self.assertIn(NativeObjectKind.IMAGE, kinds)
            self.assertIn(NativeObjectKind.FORMULA, kinds)
            self.assertIn(NativeObjectKind.CHEMICAL_REACTION, kinds)
            self.assertTrue(first_page.observed_order_signals)
            self.assertTrue(first_page.reading_order_hypotheses)
            self.assertTrue(
                all(
                    item.bbox is not None
                    for item in first_page.native_objects
                )
            )
            table = next(
                item.table
                for item in first_page.native_objects
                if item.kind is NativeObjectKind.TABLE
            )
            self.assertIsNotNone(table)
            assert table is not None
            self.assertEqual(len(table.cells), 4)
            body_cells = [cell for cell in table.cells if not cell.is_header]
            self.assertTrue(
                all(cell.header_cell_refs for cell in body_cells)
            )
            self.assertTrue(
                any(
                    item.reason_code == "empty_slide_observation"
                    and item.page_id == source.pages[1].page_id
                    for item in source.unresolved_regions
                )
            )
            self.assertEqual(source, parse_source(source_path))

    def test_pdf_retains_blank_pages_and_nested_outline(self):
        with tempfile.TemporaryDirectory() as tmp:
            source_path = Path(tmp) / "outline.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=612, height=792)
            writer.add_blank_page(width=612, height=792)
            parent = writer.add_outline_item("Chapter", 0)
            writer.add_outline_item("Section", 1, parent=parent)
            with source_path.open("wb") as handle:
                writer.write(handle)

            source = parse_source(source_path)

            self.assertEqual(len(source.pages), 2)
            self.assertEqual(
                [item.label for item in source.outline_entries],
                ["Chapter", "Section"],
            )
            self.assertEqual(
                [item.observed_level for item in source.outline_entries],
                [1, 2],
            )
            self.assertEqual(
                [item.target_page_id for item in source.outline_entries],
                [source.pages[0].page_id, source.pages[1].page_id],
            )
            self.assertEqual(len(source.unresolved_regions), 2)
            self.assertEqual(
                len({page.page_id for page in source.pages}),
                2,
            )

    def test_docx_preserves_heading_table_and_unresolved_pagination(self):
        with tempfile.TemporaryDirectory() as tmp:
            source_path = Path(tmp) / "course.docx"
            document = Document()
            document.add_heading("Aldehydes and Ketones", level=1)
            document.add_paragraph("A carbonyl group contains C=O.")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Class"
            table.cell(0, 1).text = "Position"
            table.cell(1, 0).text = "Ketone"
            table.cell(1, 1).text = "Internal"
            document.save(source_path)

            source = parse_source(source_path)

            self.assertEqual(len(source.pages), 1)
            self.assertEqual(
                [item.label for item in source.outline_entries],
                ["Aldehydes and Ketones"],
            )
            table_object = next(
                item
                for item in source.pages[0].native_objects
                if item.kind is NativeObjectKind.TABLE
            )
            assert table_object.table is not None
            self.assertEqual(len(table_object.table.cells), 4)
            self.assertTrue(
                all(
                    cell.header_cell_refs
                    for cell in table_object.table.cells
                    if not cell.is_header
                )
            )
            self.assertEqual(
                source.unresolved_regions[0].reason_code,
                "layout_pagination_unavailable",
            )

    def test_shadow_store_archives_source_before_inventory(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_path = root / "course.txt"
            source_path.write_text("Course title\nA fact.\n", encoding="utf-8")
            store = LocalArtifactStore(root / "shadow")

            result = run_source_shadow(
                source_path,
                owner_id="tenant-a",
                store=store,
            )

            self.assertEqual(
                result.source_envelope.artifact_type,
                ArtifactType.SOURCE_OBSERVATION_IR,
            )
            self.assertEqual(
                result.inventory_envelope.artifact_type,
                ArtifactType.SOURCE_INVENTORY,
            )
            self.assertEqual(
                result.inventory_envelope.input_refs,
                (store.ref(result.source_envelope),),
            )
            stored_source = store.get(
                owner_id="tenant-a",
                artifact_id=result.source_envelope.artifact_id,
            )
            stored_inventory = store.get(
                owner_id="tenant-a",
                artifact_id=result.inventory_envelope.artifact_id,
            )
            self.assertEqual(stored_source.payload, result.source_observation)
            self.assertEqual(stored_inventory.payload, result.source_inventory)


if __name__ == "__main__":
    unittest.main()
