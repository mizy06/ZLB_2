from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from backend.app.mindmap_engine.normalize import normalize_graph
from backend.app.mindmap_engine.schemas import (
    AssembleRequest,
    CropRequest,
    EvidenceRef,
    NodeCandidateIn,
    NormalizeRequest,
    NormalizedGraph,
    NormalizedNode,
    NormalizedParentCandidate,
    ParentCandidateIn,
    SolveRequest,
    VisualRegion,
)
from backend.app.mindmap_engine.service import assemble_mindmap
from backend.app.mindmap_engine.topology import solve_topology
from backend.app.mindmap_engine.visuals import crop_regions, render_document


def evidence(unit_id: str, excerpt: str) -> EvidenceRef:
    return EvidenceRef(unit_id=unit_id, excerpt=excerpt)


class MindMapEngineTests(unittest.TestCase):
    def test_normalize_merges_duplicates_and_generates_root(self):
        normalized = normalize_graph(
            NormalizeRequest(
                document_id="doc_demo",
                document_title="机器学习基础",
                nodes=[
                    NodeCandidateIn(
                        temp_id="n1",
                        name="监督学习",
                        definition="使用带标签样本学习映射。",
                        confidence=0.8,
                        evidence=[evidence("u1", "监督学习使用带标签样本")],
                    ),
                    NodeCandidateIn(
                        temp_id="n2",
                        name="监督 学习",
                        definition="从带标签数据中学习。",
                        confidence=0.7,
                        evidence=[evidence("u2", "从带标签数据中学习")],
                    ),
                ],
            )
        )

        self.assertEqual(len(normalized.nodes), 2)
        root = next(node for node in normalized.nodes if node.is_root_candidate)
        concept = next(node for node in normalized.nodes if not node.is_root_candidate)
        self.assertEqual(root.name, "机器学习基础")
        self.assertEqual(set(concept.temp_ids), {"n1", "n2"})
        self.assertTrue(normalized.parent_candidates)

    def test_assemble_builds_one_rooted_tree(self):
        request = AssembleRequest(
            document_id="doc_ml",
            document_title="机器学习",
            nodes=[
                NodeCandidateIn(
                    temp_id="root",
                    name="机器学习",
                    type="root_topic",
                    role="root_topic",
                    origin="synthesized_root",
                    is_root_candidate=True,
                    confidence=0.95,
                    support_unit_ids=["u1", "u2"],
                ),
                NodeCandidateIn(
                    temp_id="branch",
                    name="学习范式",
                    type="branch_topic",
                    role="branch_topic",
                    origin="abstractive",
                    confidence=0.9,
                    support_unit_ids=["u1", "u2"],
                ),
                NodeCandidateIn(
                    temp_id="supervised",
                    name="监督学习",
                    role="concept",
                    confidence=0.9,
                    evidence=[evidence("u1", "监督学习使用带标签样本")],
                ),
                NodeCandidateIn(
                    temp_id="unsupervised",
                    name="无监督学习",
                    role="concept",
                    confidence=0.88,
                    evidence=[evidence("u2", "无监督学习处理无标签数据")],
                ),
            ],
            parent_candidates=[
                ParentCandidateIn(
                    parent="root",
                    child="branch",
                    score=0.96,
                    classification="direct_parent",
                    verifier_score=0.95,
                ),
                ParentCandidateIn(
                    parent="branch",
                    child="supervised",
                    score=0.94,
                    classification="direct_parent",
                    verifier_score=0.94,
                ),
                ParentCandidateIn(
                    parent="branch",
                    child="unsupervised",
                    score=0.93,
                    classification="direct_parent",
                    verifier_score=0.93,
                ),
            ],
        )

        result = assemble_mindmap(request)

        self.assertTrue(result.quality.topology_valid)
        self.assertEqual(result.quality.root_count, 1)
        self.assertEqual(len(result.tree_edges), len(result.nodes) - 1)
        self.assertEqual(result.root_id, next(
            node.id for node in result.nodes if node.name == "机器学习"
        ))

    def test_provisional_edge_is_never_silent(self):
        root = NormalizedNode(
            id="root",
            temp_ids=["root"],
            name="课程",
            type="root_topic",
            role="root_topic",
            definition="",
            aliases=[],
            origin="synthesized_root",
            confidence=0.9,
            optional=False,
            activation_score=0.9,
            activation_cost=0,
            is_root_candidate=True,
            evidence=[evidence("title", "课程")],
            support_unit_ids=["u1"],
            media_asset_ids=[],
        )
        child = NormalizedNode(
            id="child",
            temp_ids=["child"],
            name="孤立知识点",
            type="concept",
            role="concept",
            definition="",
            aliases=[],
            origin="explicit",
            confidence=0.8,
            optional=False,
            activation_score=0.8,
            activation_cost=0,
            is_root_candidate=False,
            evidence=[evidence("u1", "孤立知识点")],
            support_unit_ids=[],
            media_asset_ids=[],
        )
        graph = NormalizedGraph(
            document_id="doc",
            document_title="课程",
            nodes=[root, child],
            parent_candidates=[
                NormalizedParentCandidate(
                    parent_id="root",
                    child_id="child",
                    score=0.1,
                    classification="uncertain",
                    provisional=True,
                )
            ],
            cross_links=[],
        )

        result = solve_topology(SolveRequest(graph=graph))

        self.assertTrue(result.quality.topology_valid)
        self.assertEqual(result.quality.provisional_edge_count, 1)
        self.assertTrue(any(
            item.type == "competing_parent"
            for item in result.review_items
        ))

    def test_visual_render_extracts_native_pptx_picture_and_crops_image(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_image = root / "source.png"
            Image.new("RGB", (400, 240), "white").save(source_image)

            pptx_path = root / "course.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(
                str(source_image),
                Inches(1),
                Inches(1),
                width=Inches(4),
            )
            presentation.save(pptx_path)

            rendered = render_document(
                pptx_path,
                "course.pptx",
                root / "data",
            )

            self.assertTrue(rendered.native_visuals)
            self.assertEqual(rendered.native_visuals[0].status, "ready")
            self.assertTrue(rendered.native_visuals[0].filename)

            page_render = render_document(
                source_image,
                "source.png",
                root / "data",
            )
            units = crop_regions(
                CropRequest(
                    render_id=page_render.render_id,
                    regions=[
                        VisualRegion(
                            page=1,
                            bbox=[0.25, 0.25, 0.5, 0.5],
                            visual_kind="diagram",
                            summary="中心区域",
                            knowledge_claims=["示意图中心区域"],
                        )
                    ],
                ),
                root / "data",
            )

            self.assertEqual(len(units), 1)
            self.assertGreaterEqual(units[0].asset.width or 0, 200)
            self.assertLessEqual(units[0].asset.width or 0, 210)
            self.assertGreaterEqual(units[0].asset.height or 0, 120)
            self.assertLessEqual(units[0].asset.height or 0, 130)


if __name__ == "__main__":
    unittest.main()
