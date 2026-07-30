from __future__ import annotations

import hashlib
import re
from importlib import metadata
from pathlib import Path
from typing import Any, Iterable

from docx import Document
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

from backend.vnext.artifacts.canonical import (
    canonical_json_bytes,
    payload_digest,
    stable_source_id,
)
from backend.vnext.contracts.common import (
    InterpretationStatus,
    ProducerRef,
    StringValue,
)
from backend.vnext.contracts.evidence import EvidenceNamespace, EvidenceRef
from backend.vnext.contracts.source import (
    BlockIR,
    BlockKind,
    BoundingBox,
    CharSpan,
    ChemicalReactionIR,
    CoordinateUnit,
    FormulaIR,
    HypothesisKind,
    NativeObjectIR,
    NativeObjectKind,
    ObservedOrderSignal,
    OrderSignalKind,
    OutlineEntryIR,
    PageDimensions,
    PageIR,
    PageRole,
    ParserManifest,
    ReadingOrderHypothesis,
    RoleHypothesis,
    SourceLayer,
    SourceObservationIR,
    TableCellIR,
    TableIR,
    UnresolvedRegionIR,
)


PARSER_NAME = "zlb-vnext-source-shadow"
PARSER_VERSION = "1.0.0"
PARSER_MAJOR = 1
SUPPORTED_SUFFIXES = frozenset(
    {".pdf", ".pptx", ".docx", ".txt", ".md", ".markdown"}
)

_REACTION_ARROW = re.compile(r"(?:<=>|<->|=>|->|⇌|↔|→|⇒)")
_FORMULA_RELATION = re.compile(r"[=≈≃≅≤≥<>∫∑Σ√]")
_TOC_CUE = re.compile(r"(?:目录|contents?|agenda)", re.IGNORECASE)
_REVIEW_CUE = re.compile(r"(?:复习|回顾|总结|review|summary)", re.IGNORECASE)
_EXERCISE_CUE = re.compile(
    r"(?:练习|习题|作业|exercise|quiz|questions?)",
    re.IGNORECASE,
)
_ANSWER_CUE = re.compile(r"(?:答案|解析|answer|solution)", re.IGNORECASE)
_RESEARCH_CUE = re.compile(
    r"(?:拓展|延伸|研究进展|research|further reading)",
    re.IGNORECASE,
)
_MARKDOWN_HEADING = re.compile(r"^(?P<marks>#{1,6})\s+(?P<label>.+?)\s*$")
_NUMBERED_HEADING = re.compile(
    r"^(?:第\s*[一二三四五六七八九十百\d]+\s*[章节]|"
    r"\d+(?:\.\d+){0,5})[\s:：、.-]*(?P<label>.*)$"
)

_PRODUCER = ProducerRef(
    producer_id="vnext-source-observer",
    producer_version=PARSER_VERSION,
)


def _sha256_bytes(value: bytes) -> str:
    return "sha256:" + hashlib.sha256(value).hexdigest()


def _source_hash(path: Path) -> str:
    hasher = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            hasher.update(chunk)
    return "sha256:" + hasher.hexdigest()


def _dependency_version(distribution: str) -> str:
    try:
        return metadata.version(distribution)
    except metadata.PackageNotFoundError:
        return "unavailable"


def _parser_manifest(suffix: str) -> ParserManifest:
    dependencies = {
        "python-docx": _dependency_version("python-docx"),
        "python-pptx": _dependency_version("python-pptx"),
        "pypdf": _dependency_version("pypdf"),
    }
    return ParserManifest(
        parser_name=PARSER_NAME,
        parser_version=PARSER_VERSION,
        parser_major=PARSER_MAJOR,
        dependency_digest=payload_digest(dependencies),
        configuration=(
            StringValue(key="source_type", value=suffix.lstrip(".")),
            StringValue(
                key="render_reference_mode",
                value="logical-source-page",
            ),
        ),
    )


def _stable_hypothesis_id(
    *,
    source_hash: str,
    locator: Any,
) -> str:
    identity = {
        "locator": locator,
        "parser_major": PARSER_MAJOR,
        "source_hash": source_hash,
    }
    digest = hashlib.sha256(
        b"zlb-vnext-hypothesis-id-v1\0"
        + canonical_json_bytes(identity)
    ).hexdigest()
    return "hyp_" + digest[:32]


def _courseware_evidence(
    ref_id: str,
    *,
    content: bytes | str | None = None,
) -> EvidenceRef:
    digest: str | None = None
    if isinstance(content, str):
        digest = _sha256_bytes(content.encode("utf-8"))
    elif isinstance(content, bytes):
        digest = _sha256_bytes(content)
    return EvidenceRef(
        namespace=EvidenceNamespace.COURSEWARE,
        ref_id=ref_id,
        content_digest=digest,
    )


def _source_id(
    kind: str,
    *,
    source_hash: str,
    locator: Any,
) -> str:
    return stable_source_id(
        kind,
        source_hash=source_hash,
        parser_major=PARSER_MAJOR,
        locator=locator,
    )


def _page_dimensions(width: float, height: float) -> PageDimensions:
    return PageDimensions(
        width=max(float(width), 1.0),
        height=max(float(height), 1.0),
        unit=CoordinateUnit.POINT,
    )


def _full_page_bbox(dimensions: PageDimensions) -> BoundingBox:
    return BoundingBox(
        x=0,
        y=0,
        width=dimensions.width,
        height=dimensions.height,
        unit=dimensions.unit,
    )


def _emu_to_points(value: Any) -> float:
    try:
        return max(float(value) / 12700.0, 0.0)
    except (TypeError, ValueError):
        return 0.0


def _shape_bbox(shape: Any) -> BoundingBox | None:
    width = _emu_to_points(getattr(shape, "width", 0))
    height = _emu_to_points(getattr(shape, "height", 0))
    if width <= 0 or height <= 0:
        return None
    return BoundingBox(
        x=_emu_to_points(getattr(shape, "left", 0)),
        y=_emu_to_points(getattr(shape, "top", 0)),
        width=width,
        height=height,
        unit=CoordinateUnit.POINT,
    )


def _outline_level(label: str, default: int = 1) -> int:
    stripped = label.strip()
    match = re.match(r"^(?P<number>\d+(?:\.\d+)*)", stripped)
    if match:
        return min(match.group("number").count(".") + 1, 6)
    if _NUMBERED_HEADING.match(stripped):
        return default
    return default


def _classify_text_object(
    text: str,
    *,
    page_render_ref: EvidenceRef,
    xml_math: bool = False,
) -> tuple[
    NativeObjectKind,
    FormulaIR | None,
    ChemicalReactionIR | None,
    InterpretationStatus,
]:
    body = text.strip()
    arrow = _REACTION_ARROW.search(body)
    if arrow and re.search(r"[\w\u3400-\u9fff]", body):
        return (
            NativeObjectKind.CHEMICAL_REACTION,
            None,
            ChemicalReactionIR(
                reactants=(),
                products=(),
                arrow=arrow.group(0),
                parse_status="unresolved",
            ),
            InterpretationStatus.INFERRED,
        )
    if xml_math or (
        _FORMULA_RELATION.search(body)
        and re.search(r"[A-Za-zΑ-ω\d]", body)
    ):
        return (
            NativeObjectKind.FORMULA,
            FormulaIR(
                display_text=body,
                render_ref=page_render_ref,
                parse_status="unresolved",
            ),
            None,
            (
                InterpretationStatus.OBSERVED
                if xml_math
                else InterpretationStatus.INFERRED
            ),
        )
    return (
        NativeObjectKind.TEXT,
        None,
        None,
        InterpretationStatus.OBSERVED,
    )


def _role_for_page(
    *,
    physical_index: int,
    text: str,
    title: str | None,
    content_block_count: int,
) -> tuple[PageRole, float]:
    combined = "\n".join(part for part in (title, text) if part).strip()
    if not combined:
        return PageRole.UNKNOWN, 0.25
    if _TOC_CUE.search(combined):
        return PageRole.TOC, 0.9
    if _ANSWER_CUE.search(combined):
        return PageRole.ANSWER, 0.8
    if _EXERCISE_CUE.search(combined):
        return PageRole.EXERCISE, 0.8
    if _REVIEW_CUE.search(combined):
        return PageRole.REVIEW, 0.75
    if _RESEARCH_CUE.search(combined):
        return PageRole.RESEARCH_ASIDE, 0.7
    if physical_index == 0 and title and content_block_count <= 2:
        return PageRole.COVER, 0.75
    if title and content_block_count <= 1:
        return PageRole.SECTION_DIVIDER, 0.65
    return PageRole.CONTENT, 0.7


def _role_hypothesis(
    *,
    source_hash: str,
    page_id: str,
    physical_index: int,
    text: str,
    title: str | None,
    content_block_count: int,
) -> RoleHypothesis:
    role, confidence = _role_for_page(
        physical_index=physical_index,
        text=text,
        title=title,
        content_block_count=content_block_count,
    )
    return RoleHypothesis(
        hypothesis_id=_stable_hypothesis_id(
            source_hash=source_hash,
            locator={
                "kind": "page_role",
                "page": physical_index,
                "role": role.value,
            },
        ),
        role=role,
        confidence=confidence,
        interpretation_status=InterpretationStatus.INFERRED,
        producer=_PRODUCER,
        evidence_refs=(_courseware_evidence(page_id),),
    )


def _order_records(
    *,
    source_hash: str,
    physical_index: int,
    page_id: str,
    observed_ids: tuple[str, ...],
    ordered_blocks: tuple[BlockIR, ...],
) -> tuple[
    tuple[ObservedOrderSignal, ...],
    tuple[ReadingOrderHypothesis, ...],
]:
    signals: tuple[ObservedOrderSignal, ...] = ()
    if len(observed_ids) >= 2:
        signal_id = _source_id(
            "order",
            source_hash=source_hash,
            locator={
                "page": physical_index,
                "kind": "native_sequence",
            },
        )
        signals = (
            ObservedOrderSignal(
                signal_id=signal_id,
                kind=OrderSignalKind.NATIVE_SEQUENCE,
                ordered_source_ids=observed_ids,
                producer=_PRODUCER,
                evidence_refs=(_courseware_evidence(page_id),),
            ),
        )
    hypotheses: tuple[ReadingOrderHypothesis, ...] = ()
    if len(ordered_blocks) >= 2:
        hypotheses = (
            ReadingOrderHypothesis(
                hypothesis_id=_stable_hypothesis_id(
                    source_hash=source_hash,
                    locator={
                        "kind": "reading_order",
                        "page": physical_index,
                    },
                ),
                ordered_block_ids=tuple(
                    block.block_id for block in ordered_blocks
                ),
                confidence=0.75,
                interpretation_status=InterpretationStatus.INFERRED,
                producer=_PRODUCER,
                evidence_refs=(_courseware_evidence(page_id),),
            ),
        )
    return signals, hypotheses


def _block_sort_key(block: BlockIR) -> tuple[float, float, int]:
    if block.bbox is None:
        return (float(block.native_order_hint or 0), 0.0, 0)
    return (
        block.bbox.y,
        block.bbox.x,
        block.native_order_hint or 0,
    )


def _table_cells(
    *,
    source_hash: str,
    table_locator: Any,
    rows: list[list[tuple[str, int, int]]],
) -> TableIR:
    row_count = max(len(rows), 1)
    column_count = max(
        (sum(item[2] for item in row) for row in rows),
        default=1,
    )
    raw_cells: list[dict[str, Any]] = []
    header_by_column: dict[int, str] = {}
    for row_index, row in enumerate(rows):
        column_index = 0
        for text, row_span, column_span in row:
            cell_id = _source_id(
                "cell",
                source_hash=source_hash,
                locator={
                    "table": table_locator,
                    "row": row_index,
                    "column": column_index,
                },
            )
            item = {
                "cell_id": cell_id,
                "row_index": row_index,
                "column_index": column_index,
                "row_span": max(row_span, 1),
                "column_span": max(column_span, 1),
                "text": text,
                "is_header": row_index == 0,
            }
            raw_cells.append(item)
            if row_index == 0:
                for column in range(
                    column_index,
                    column_index + item["column_span"],
                ):
                    header_by_column[column] = cell_id
            column_index += item["column_span"]

    cells = tuple(
        TableCellIR(
            **item,
            header_cell_refs=(
                ()
                if item["is_header"]
                else tuple(
                    dict.fromkeys(
                        header_by_column[column]
                        for column in range(
                            item["column_index"],
                            item["column_index"]
                            + item["column_span"],
                        )
                        if column in header_by_column
                    )
                )
            ),
            evidence_refs=(
                _courseware_evidence(
                    item["cell_id"],
                    content=item["text"],
                ),
            ),
        )
        for item in raw_cells
    )
    return TableIR(
        row_count=row_count,
        column_count=column_count,
        cells=cells,
    )


def _pptx_table_rows(table: Any) -> list[list[tuple[str, int, int]]]:
    rows: list[list[tuple[str, int, int]]] = []
    for row in table.rows:
        values: list[tuple[str, int, int]] = []
        cells = list(row.cells)
        index = 0
        while index < len(cells):
            cell = cells[index]
            run = 1
            while (
                index + run < len(cells)
                and cells[index + run]._tc is cell._tc
            ):
                run += 1
            span_width = int(getattr(cell, "span_width", run) or run)
            span_height = int(getattr(cell, "span_height", 1) or 1)
            values.append(
                (
                    cell.text.strip(),
                    max(span_height, 1),
                    max(span_width, run, 1),
                )
            )
            index += run
        rows.append(values)
    return rows


def _docx_table_rows(table: DocxTable) -> list[list[tuple[str, int, int]]]:
    rows: list[list[tuple[str, int, int]]] = []
    for row in table.rows:
        values: list[tuple[str, int, int]] = []
        cells = list(row.cells)
        index = 0
        while index < len(cells):
            cell = cells[index]
            run = 1
            while (
                index + run < len(cells)
                and cells[index + run]._tc is cell._tc
            ):
                run += 1
            values.append((cell.text.strip(), 1, run))
            index += run
        rows.append(values)
    return rows


def _table_text(rows: list[list[tuple[str, int, int]]]) -> str:
    return "\n".join(
        " | ".join(text for text, _, _ in row).strip()
        for row in rows
        if any(text for text, _, _ in row)
    ).strip()


def _shape_has_math(shape: Any) -> bool:
    try:
        return bool(
            shape.element.xpath(
                ".//*[local-name()='oMath' or local-name()='oMathPara']"
            )
        )
    except (AttributeError, TypeError, ValueError):
        return False


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return str(getattr(shape, "text", "")).strip()
    if getattr(shape, "has_chart", False):
        chart = shape.chart
        if chart.has_title:
            return chart.chart_title.text_frame.text.strip()
    return ""


def _pptx_shape_kind(shape: Any) -> NativeObjectKind:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return NativeObjectKind.IMAGE
    if shape_type == MSO_SHAPE_TYPE.TABLE or getattr(
        shape,
        "has_table",
        False,
    ):
        return NativeObjectKind.TABLE
    if shape_type == MSO_SHAPE_TYPE.CHART or getattr(
        shape,
        "has_chart",
        False,
    ):
        return NativeObjectKind.CHART
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return NativeObjectKind.GROUP
    if getattr(shape, "has_text_frame", False):
        return NativeObjectKind.TEXT
    return NativeObjectKind.SHAPE


def _parse_pptx(
    path: Path,
    *,
    source_hash: str,
) -> tuple[
    tuple[PageIR, ...],
    tuple[OutlineEntryIR, ...],
    tuple[UnresolvedRegionIR, ...],
]:
    presentation = Presentation(str(path))
    dimensions = _page_dimensions(
        _emu_to_points(presentation.slide_width),
        _emu_to_points(presentation.slide_height),
    )
    pages: list[PageIR] = []
    outline: list[OutlineEntryIR] = []
    unresolved: list[UnresolvedRegionIR] = []
    outline_order = 0

    for physical_index, slide in enumerate(presentation.slides):
        page_id = _source_id(
            "page",
            source_hash=source_hash,
            locator={"slide": physical_index},
        )
        render_id = _source_id(
            "render",
            source_hash=source_hash,
            locator={"slide": physical_index, "kind": "full_slide"},
        )
        render_ref = _courseware_evidence(
            render_id,
            content=source_hash,
        )
        objects: list[NativeObjectIR] = []
        blocks: list[BlockIR] = []
        native_sequence: list[str] = []
        order_counter = 0

        def visit_shape(
            shape: Any,
            *,
            shape_path: tuple[int, ...],
            parent_group_id: str | None,
        ) -> None:
            nonlocal order_counter
            locator = {
                "slide": physical_index,
                "shape_path": shape_path,
            }
            object_id = _source_id(
                "object",
                source_hash=source_hash,
                locator=locator,
            )
            kind = _pptx_shape_kind(shape)
            text = _shape_text(shape)
            formula: FormulaIR | None = None
            reaction: ChemicalReactionIR | None = None
            status = InterpretationStatus.OBSERVED
            if kind is NativeObjectKind.TEXT and text:
                kind, formula, reaction, status = _classify_text_object(
                    text,
                    page_render_ref=render_ref,
                    xml_math=_shape_has_math(shape),
                )

            asset_ref: EvidenceRef | None = None
            evidence_content: bytes | str | None = text
            if kind is NativeObjectKind.IMAGE:
                try:
                    image_blob = shape.image.blob
                except (AttributeError, KeyError, ValueError):
                    image_blob = b""
                asset_id = _source_id(
                    "asset",
                    source_hash=source_hash,
                    locator={**locator, "kind": "image_blob"},
                )
                asset_ref = _courseware_evidence(
                    asset_id,
                    content=image_blob or None,
                )
                evidence_content = image_blob or None

            table_payload: TableIR | None = None
            table_rows: list[list[tuple[str, int, int]]] = []
            if kind is NativeObjectKind.TABLE:
                table_rows = _pptx_table_rows(shape.table)
                table_payload = _table_cells(
                    source_hash=source_hash,
                    table_locator=locator,
                    rows=table_rows,
                )
                text = _table_text(table_rows)
                evidence_content = text

            evidence = _courseware_evidence(
                object_id,
                content=evidence_content,
            )
            objects.append(
                NativeObjectIR(
                    object_id=object_id,
                    page_id=page_id,
                    kind=kind,
                    bbox=_shape_bbox(shape),
                    native_order_hint=order_counter,
                    parent_group_id=parent_group_id,
                    text=text,
                    asset_ref=asset_ref,
                    table=table_payload,
                    formula=formula,
                    reaction=reaction,
                    source_layer=SourceLayer.NATIVE,
                    observation_status=status,
                    producer=_PRODUCER,
                    evidence_refs=(evidence,),
                )
            )
            native_sequence.append(object_id)

            if text:
                block_id = _source_id(
                    "block",
                    source_hash=source_hash,
                    locator={**locator, "projection": "text"},
                )
                is_title = (
                    slide.shapes.title is not None
                    and getattr(shape, "shape_id", None)
                    == slide.shapes.title.shape_id
                )
                block_kind = (
                    BlockKind.TITLE
                    if is_title
                    else (
                        BlockKind.TABLE_TEXT
                        if kind is NativeObjectKind.TABLE
                        else (
                            BlockKind.FORMULA_TEXT
                            if kind is NativeObjectKind.FORMULA
                            else (
                                BlockKind.REACTION_TEXT
                                if kind
                                is NativeObjectKind.CHEMICAL_REACTION
                                else BlockKind.PARAGRAPH
                            )
                        )
                    )
                )
                blocks.append(
                    BlockIR(
                        block_id=block_id,
                        page_id=page_id,
                        kind=block_kind,
                        text=text,
                        bbox=_shape_bbox(shape),
                        native_order_hint=order_counter,
                        native_object_id=object_id,
                        parent_group_id=parent_group_id,
                        source_layer=SourceLayer.NATIVE,
                        source_confidence=1.0,
                        char_spans=(
                            CharSpan(start=0, end=len(text), text=text),
                        ),
                        observation_status=InterpretationStatus.OBSERVED,
                        producer=_PRODUCER,
                        evidence_refs=(
                            _courseware_evidence(
                                block_id,
                                content=text,
                            ),
                        ),
                    )
                )
            order_counter += 1

            if kind is NativeObjectKind.GROUP:
                for child_index, child in enumerate(shape.shapes):
                    visit_shape(
                        child,
                        shape_path=(*shape_path, child_index),
                        parent_group_id=object_id,
                    )

        for shape_index, shape in enumerate(slide.shapes):
            visit_shape(
                shape,
                shape_path=(shape_index,),
                parent_group_id=None,
            )

        title = None
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip() or None
        if title:
            outline_id = _source_id(
                "outline",
                source_hash=source_hash,
                locator={"slide": physical_index, "title": True},
            )
            outline.append(
                OutlineEntryIR(
                    outline_entry_id=outline_id,
                    label=title,
                    observed_level=_outline_level(
                        title,
                        default=1 if physical_index == 0 else 2,
                    ),
                    source_order=outline_order,
                    target_page_id=page_id,
                    native_target=f"slide:{physical_index + 1}",
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(
                            outline_id,
                            content=title,
                        ),
                    ),
                )
            )
            outline_order += 1

        if not blocks and not objects:
            region_id = _source_id(
                "unresolved",
                source_hash=source_hash,
                locator={
                    "slide": physical_index,
                    "reason": "empty_slide_observation",
                },
            )
            unresolved.append(
                UnresolvedRegionIR(
                    region_source_id=region_id,
                    page_id=page_id,
                    bbox=_full_page_bbox(dimensions),
                    reason_code="empty_slide_observation",
                    producer=_PRODUCER,
                    evidence_refs=(render_ref,),
                )
            )

        ordered_blocks = tuple(sorted(blocks, key=_block_sort_key))
        signals, reading = _order_records(
            source_hash=source_hash,
            physical_index=physical_index,
            page_id=page_id,
            observed_ids=tuple(native_sequence),
            ordered_blocks=ordered_blocks,
        )
        page_text = "\n".join(block.text for block in ordered_blocks)
        pages.append(
            PageIR(
                page_id=page_id,
                physical_index=physical_index,
                logical_number=str(physical_index + 1),
                dimensions=dimensions,
                blocks=tuple(blocks),
                native_objects=tuple(objects),
                observed_order_signals=signals,
                reading_order_hypotheses=reading,
                role_hypotheses=(
                    _role_hypothesis(
                        source_hash=source_hash,
                        page_id=page_id,
                        physical_index=physical_index,
                        text=page_text,
                        title=title,
                        content_block_count=len(blocks),
                    ),
                ),
                render_ref=render_ref,
            )
        )
    return tuple(pages), tuple(outline), tuple(unresolved)


def _pdf_outline_entries(
    reader: PdfReader,
    *,
    source_hash: str,
    page_ids: tuple[str, ...],
) -> tuple[OutlineEntryIR, ...]:
    entries: list[OutlineEntryIR] = []
    source_order = 0

    def walk(items: Iterable[Any], level: int) -> None:
        nonlocal source_order
        for item in items:
            if isinstance(item, list):
                walk(item, level + 1)
                continue
            label = str(getattr(item, "title", "") or "").strip()
            if not label:
                continue
            target_page_id: str | None = None
            native_target: str | None = None
            try:
                page_index = reader.get_destination_page_number(item)
            except (AttributeError, KeyError, TypeError, ValueError):
                page_index = None
            if (
                isinstance(page_index, int)
                and 0 <= page_index < len(page_ids)
            ):
                target_page_id = page_ids[page_index]
                native_target = f"page:{page_index + 1}"
            outline_id = _source_id(
                "outline",
                source_hash=source_hash,
                locator={
                    "pdf_outline_order": source_order,
                    "level": level,
                },
            )
            entries.append(
                OutlineEntryIR(
                    outline_entry_id=outline_id,
                    label=label,
                    observed_level=max(level, 1),
                    source_order=source_order,
                    target_page_id=target_page_id,
                    native_target=native_target,
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(
                            outline_id,
                            content=label,
                        ),
                    ),
                )
            )
            source_order += 1

    try:
        outline = reader.outline
    except (AttributeError, KeyError, TypeError, ValueError):
        outline = ()
    walk(outline, 1)
    return tuple(entries)


def _parse_pdf(
    path: Path,
    *,
    source_hash: str,
) -> tuple[
    tuple[PageIR, ...],
    tuple[OutlineEntryIR, ...],
    tuple[UnresolvedRegionIR, ...],
]:
    reader = PdfReader(str(path))
    pages: list[PageIR] = []
    unresolved: list[UnresolvedRegionIR] = []
    page_ids = tuple(
        _source_id(
            "page",
            source_hash=source_hash,
            locator={"pdf_page": index},
        )
        for index in range(len(reader.pages))
    )
    try:
        labels = tuple(reader.page_labels)
    except (AttributeError, KeyError, TypeError, ValueError):
        labels = ()

    for physical_index, page in enumerate(reader.pages):
        page_id = page_ids[physical_index]
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        dimensions = _page_dimensions(width, height)
        render_id = _source_id(
            "render",
            source_hash=source_hash,
            locator={"pdf_page": physical_index, "kind": "full_page"},
        )
        render_ref = _courseware_evidence(
            render_id,
            content=source_hash,
        )
        extraction_error: str | None = None
        try:
            text = page.extract_text(extraction_mode="layout") or ""
        except (KeyError, TypeError, ValueError) as exc:
            extraction_error = type(exc).__name__
            try:
                text = page.extract_text() or ""
            except (KeyError, TypeError, ValueError) as fallback_exc:
                text = ""
                extraction_error = type(fallback_exc).__name__
        lines = tuple(
            line.strip()
            for line in text.replace("\r\n", "\n").splitlines()
            if line.strip()
        )
        objects: list[NativeObjectIR] = []
        blocks: list[BlockIR] = []
        native_sequence: list[str] = []
        for line_index, line in enumerate(lines):
            locator = {
                "pdf_page": physical_index,
                "text_line": line_index,
            }
            object_id = _source_id(
                "object",
                source_hash=source_hash,
                locator=locator,
            )
            kind, formula, reaction, status = _classify_text_object(
                line,
                page_render_ref=render_ref,
            )
            objects.append(
                NativeObjectIR(
                    object_id=object_id,
                    page_id=page_id,
                    kind=kind,
                    native_order_hint=line_index,
                    text=line,
                    formula=formula,
                    reaction=reaction,
                    source_layer=SourceLayer.NATIVE,
                    observation_status=status,
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(
                            object_id,
                            content=line,
                        ),
                    ),
                )
            )
            native_sequence.append(object_id)
            block_id = _source_id(
                "block",
                source_hash=source_hash,
                locator={**locator, "projection": "text"},
            )
            block_kind = (
                BlockKind.FORMULA_TEXT
                if kind is NativeObjectKind.FORMULA
                else (
                    BlockKind.REACTION_TEXT
                    if kind is NativeObjectKind.CHEMICAL_REACTION
                    else BlockKind.PARAGRAPH
                )
            )
            blocks.append(
                BlockIR(
                    block_id=block_id,
                    page_id=page_id,
                    kind=block_kind,
                    text=line,
                    native_order_hint=line_index,
                    native_object_id=object_id,
                    source_layer=SourceLayer.NATIVE,
                    source_confidence=0.8,
                    char_spans=(
                        CharSpan(start=0, end=len(line), text=line),
                    ),
                    observation_status=InterpretationStatus.OBSERVED,
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(
                            block_id,
                            content=line,
                        ),
                    ),
                )
            )

        try:
            images = tuple(page.images)
        except (AttributeError, KeyError, TypeError, ValueError):
            images = ()
        for image_index, image in enumerate(images):
            image_name = str(getattr(image, "name", image_index))
            image_data = bytes(getattr(image, "data", b""))
            locator = {
                "pdf_page": physical_index,
                "image": image_index,
                "name": image_name,
            }
            object_id = _source_id(
                "object",
                source_hash=source_hash,
                locator=locator,
            )
            asset_id = _source_id(
                "asset",
                source_hash=source_hash,
                locator={**locator, "kind": "image_blob"},
            )
            objects.append(
                NativeObjectIR(
                    object_id=object_id,
                    page_id=page_id,
                    kind=NativeObjectKind.IMAGE,
                    native_order_hint=len(native_sequence),
                    asset_ref=_courseware_evidence(
                        asset_id,
                        content=image_data or None,
                    ),
                    source_layer=SourceLayer.NATIVE,
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(
                            object_id,
                            content=image_data or image_name,
                        ),
                    ),
                )
            )
            native_sequence.append(object_id)

        if not lines:
            region_id = _source_id(
                "unresolved",
                source_hash=source_hash,
                locator={
                    "pdf_page": physical_index,
                    "reason": "page_text_unresolved",
                },
            )
            unresolved.append(
                UnresolvedRegionIR(
                    region_source_id=region_id,
                    page_id=page_id,
                    bbox=_full_page_bbox(dimensions),
                    reason_code=(
                        "pdf_text_extraction_failed"
                        if extraction_error
                        else "page_text_unresolved"
                    ),
                    producer=_PRODUCER,
                    evidence_refs=(render_ref,),
                )
            )

        ordered_blocks = tuple(blocks)
        signals, reading = _order_records(
            source_hash=source_hash,
            physical_index=physical_index,
            page_id=page_id,
            observed_ids=tuple(native_sequence),
            ordered_blocks=ordered_blocks,
        )
        page_text = "\n".join(lines)
        pages.append(
            PageIR(
                page_id=page_id,
                physical_index=physical_index,
                logical_number=(
                    labels[physical_index]
                    if physical_index < len(labels)
                    else str(physical_index + 1)
                ),
                dimensions=dimensions,
                blocks=ordered_blocks,
                native_objects=tuple(objects),
                observed_order_signals=signals,
                reading_order_hypotheses=reading,
                role_hypotheses=(
                    _role_hypothesis(
                        source_hash=source_hash,
                        page_id=page_id,
                        physical_index=physical_index,
                        text=page_text,
                        title=lines[0] if lines else None,
                        content_block_count=len(blocks),
                    ),
                ),
                render_ref=render_ref,
            )
        )
    return (
        tuple(pages),
        _pdf_outline_entries(
            reader,
            source_hash=source_hash,
            page_ids=page_ids,
        ),
        tuple(unresolved),
    )


def _docx_heading_level(paragraph: Paragraph) -> int | None:
    style = getattr(paragraph, "style", None)
    name = str(getattr(style, "name", "") or "")
    match = re.match(r"Heading\s+(?P<level>[1-9])", name, re.IGNORECASE)
    if match:
        return int(match.group("level"))
    return None


def _parse_docx(
    path: Path,
    *,
    source_hash: str,
) -> tuple[
    tuple[PageIR, ...],
    tuple[OutlineEntryIR, ...],
    tuple[UnresolvedRegionIR, ...],
]:
    document = Document(str(path))
    section = document.sections[0]
    dimensions = _page_dimensions(
        _emu_to_points(section.page_width),
        _emu_to_points(section.page_height),
    )
    page_id = _source_id(
        "page",
        source_hash=source_hash,
        locator={"docx_flow": 0},
    )
    render_id = _source_id(
        "render",
        source_hash=source_hash,
        locator={"docx_flow": 0, "kind": "logical_page"},
    )
    render_ref = _courseware_evidence(render_id, content=source_hash)
    objects: list[NativeObjectIR] = []
    blocks: list[BlockIR] = []
    outline: list[OutlineEntryIR] = []
    native_sequence: list[str] = []
    title: str | None = None

    for item_index, item in enumerate(document.iter_inner_content()):
        locator = {"docx_flow": 0, "body_item": item_index}
        object_id = _source_id(
            "object",
            source_hash=source_hash,
            locator=locator,
        )
        if isinstance(item, Paragraph):
            text = item.text.strip()
            if not text:
                continue
            heading_level = _docx_heading_level(item)
            kind, formula, reaction, status = _classify_text_object(
                text,
                page_render_ref=render_ref,
            )
            evidence = _courseware_evidence(object_id, content=text)
            objects.append(
                NativeObjectIR(
                    object_id=object_id,
                    page_id=page_id,
                    kind=kind,
                    native_order_hint=item_index,
                    text=text,
                    formula=formula,
                    reaction=reaction,
                    source_layer=SourceLayer.NATIVE,
                    observation_status=status,
                    producer=_PRODUCER,
                    evidence_refs=(evidence,),
                )
            )
            native_sequence.append(object_id)
            block_id = _source_id(
                "block",
                source_hash=source_hash,
                locator={**locator, "projection": "text"},
            )
            style_name = str(
                getattr(getattr(item, "style", None), "name", "") or ""
            )
            block_kind = (
                BlockKind.HEADING
                if heading_level is not None
                else (
                    BlockKind.LIST
                    if "list" in style_name.casefold()
                    else BlockKind.PARAGRAPH
                )
            )
            blocks.append(
                BlockIR(
                    block_id=block_id,
                    page_id=page_id,
                    kind=block_kind,
                    text=text,
                    native_order_hint=item_index,
                    native_object_id=object_id,
                    source_layer=SourceLayer.NATIVE,
                    source_confidence=1.0,
                    char_spans=(
                        CharSpan(start=0, end=len(text), text=text),
                    ),
                    observation_status=InterpretationStatus.OBSERVED,
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(block_id, content=text),
                    ),
                )
            )
            if heading_level is not None:
                if title is None:
                    title = text
                outline_id = _source_id(
                    "outline",
                    source_hash=source_hash,
                    locator={**locator, "heading": heading_level},
                )
                outline.append(
                    OutlineEntryIR(
                        outline_entry_id=outline_id,
                        label=text,
                        observed_level=heading_level,
                        source_order=len(outline),
                        target_page_id=page_id,
                        native_target=f"body-item:{item_index}",
                        producer=_PRODUCER,
                        evidence_refs=(
                            _courseware_evidence(
                                outline_id,
                                content=text,
                            ),
                        ),
                    )
                )
            continue

        if isinstance(item, DocxTable):
            rows = _docx_table_rows(item)
            text = _table_text(rows)
            table = _table_cells(
                source_hash=source_hash,
                table_locator=locator,
                rows=rows,
            )
            objects.append(
                NativeObjectIR(
                    object_id=object_id,
                    page_id=page_id,
                    kind=NativeObjectKind.TABLE,
                    native_order_hint=item_index,
                    text=text,
                    table=table,
                    source_layer=SourceLayer.NATIVE,
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(
                            object_id,
                            content=text,
                        ),
                    ),
                )
            )
            native_sequence.append(object_id)
            if text:
                block_id = _source_id(
                    "block",
                    source_hash=source_hash,
                    locator={**locator, "projection": "table_text"},
                )
                blocks.append(
                    BlockIR(
                        block_id=block_id,
                        page_id=page_id,
                        kind=BlockKind.TABLE_TEXT,
                        text=text,
                        native_order_hint=item_index,
                        native_object_id=object_id,
                        source_layer=SourceLayer.NATIVE,
                        source_confidence=1.0,
                        char_spans=(
                            CharSpan(start=0, end=len(text), text=text),
                        ),
                        observation_status=InterpretationStatus.OBSERVED,
                        producer=_PRODUCER,
                        evidence_refs=(
                            _courseware_evidence(
                                block_id,
                                content=text,
                            ),
                        ),
                    )
                )

    image_index = 0
    for relationship in document.part.rels.values():
        if not relationship.reltype.endswith("/image"):
            continue
        locator = {"docx_image_relationship": relationship.rId}
        object_id = _source_id(
            "object",
            source_hash=source_hash,
            locator=locator,
        )
        asset_id = _source_id(
            "asset",
            source_hash=source_hash,
            locator={**locator, "kind": "image_blob"},
        )
        try:
            blob = relationship.target_part.blob
        except (AttributeError, KeyError, ValueError):
            blob = b""
        objects.append(
            NativeObjectIR(
                object_id=object_id,
                page_id=page_id,
                kind=NativeObjectKind.IMAGE,
                native_order_hint=len(native_sequence) + image_index,
                asset_ref=_courseware_evidence(
                    asset_id,
                    content=blob or None,
                ),
                source_layer=SourceLayer.NATIVE,
                producer=_PRODUCER,
                evidence_refs=(
                    _courseware_evidence(
                        object_id,
                        content=blob or relationship.rId,
                    ),
                ),
            )
        )
        native_sequence.append(object_id)
        image_index += 1

    ordered_blocks = tuple(blocks)
    signals, reading = _order_records(
        source_hash=source_hash,
        physical_index=0,
        page_id=page_id,
        observed_ids=tuple(native_sequence),
        ordered_blocks=ordered_blocks,
    )
    unresolved_id = _source_id(
        "unresolved",
        source_hash=source_hash,
        locator={
            "docx_flow": 0,
            "reason": "layout_pagination_unavailable",
        },
    )
    unresolved = (
        UnresolvedRegionIR(
            region_source_id=unresolved_id,
            page_id=page_id,
            bbox=_full_page_bbox(dimensions),
            reason_code="layout_pagination_unavailable",
            transcription="\n".join(block.text for block in ordered_blocks),
            producer=_PRODUCER,
            evidence_refs=(render_ref,),
        ),
    )
    page = PageIR(
        page_id=page_id,
        physical_index=0,
        logical_number=None,
        dimensions=dimensions,
        blocks=ordered_blocks,
        native_objects=tuple(objects),
        observed_order_signals=signals,
        reading_order_hypotheses=reading,
        role_hypotheses=(
            _role_hypothesis(
                source_hash=source_hash,
                page_id=page_id,
                physical_index=0,
                text="\n".join(block.text for block in ordered_blocks),
                title=title,
                content_block_count=len(blocks),
            ),
        ),
        render_ref=render_ref,
    )
    return (page,), tuple(outline), unresolved


def _parse_text(
    path: Path,
    *,
    source_hash: str,
) -> tuple[
    tuple[PageIR, ...],
    tuple[OutlineEntryIR, ...],
    tuple[UnresolvedRegionIR, ...],
]:
    text = path.read_text(encoding="utf-8", errors="replace")
    dimensions = PageDimensions(
        width=1,
        height=1,
        unit=CoordinateUnit.NORMALIZED,
    )
    page_id = _source_id(
        "page",
        source_hash=source_hash,
        locator={"text_flow": 0},
    )
    render_id = _source_id(
        "render",
        source_hash=source_hash,
        locator={"text_flow": 0, "kind": "logical_page"},
    )
    render_ref = _courseware_evidence(render_id, content=text)
    blocks: list[BlockIR] = []
    objects: list[NativeObjectIR] = []
    outline: list[OutlineEntryIR] = []
    native_sequence: list[str] = []
    title: str | None = None

    for line_index, raw_line in enumerate(text.splitlines()):
        line = raw_line.strip()
        if not line:
            continue
        heading = _MARKDOWN_HEADING.match(line)
        label = heading.group("label").strip() if heading else line
        locator = {"text_flow": 0, "line": line_index}
        object_id = _source_id(
            "object",
            source_hash=source_hash,
            locator=locator,
        )
        kind, formula, reaction, status = _classify_text_object(
            line,
            page_render_ref=render_ref,
        )
        objects.append(
            NativeObjectIR(
                object_id=object_id,
                page_id=page_id,
                kind=kind,
                native_order_hint=line_index,
                text=line,
                formula=formula,
                reaction=reaction,
                source_layer=SourceLayer.NATIVE,
                observation_status=status,
                producer=_PRODUCER,
                evidence_refs=(
                    _courseware_evidence(object_id, content=line),
                ),
            )
        )
        native_sequence.append(object_id)
        block_id = _source_id(
            "block",
            source_hash=source_hash,
            locator={**locator, "projection": "text"},
        )
        blocks.append(
            BlockIR(
                block_id=block_id,
                page_id=page_id,
                kind=(
                    BlockKind.HEADING
                    if heading
                    else BlockKind.PARAGRAPH
                ),
                text=line,
                native_order_hint=line_index,
                native_object_id=object_id,
                source_layer=SourceLayer.NATIVE,
                source_confidence=1.0,
                char_spans=(CharSpan(start=0, end=len(line), text=line),),
                observation_status=InterpretationStatus.OBSERVED,
                producer=_PRODUCER,
                evidence_refs=(
                    _courseware_evidence(block_id, content=line),
                ),
            )
        )
        if heading:
            if title is None:
                title = label
            outline_id = _source_id(
                "outline",
                source_hash=source_hash,
                locator={**locator, "heading": True},
            )
            outline.append(
                OutlineEntryIR(
                    outline_entry_id=outline_id,
                    label=label,
                    observed_level=len(heading.group("marks")),
                    source_order=len(outline),
                    target_page_id=page_id,
                    native_target=f"line:{line_index + 1}",
                    producer=_PRODUCER,
                    evidence_refs=(
                        _courseware_evidence(
                            outline_id,
                            content=label,
                        ),
                    ),
                )
            )

    ordered_blocks = tuple(blocks)
    signals, reading = _order_records(
        source_hash=source_hash,
        physical_index=0,
        page_id=page_id,
        observed_ids=tuple(native_sequence),
        ordered_blocks=ordered_blocks,
    )
    unresolved: tuple[UnresolvedRegionIR, ...] = ()
    if not blocks:
        unresolved_id = _source_id(
            "unresolved",
            source_hash=source_hash,
            locator={
                "text_flow": 0,
                "reason": "empty_text_document",
            },
        )
        unresolved = (
            UnresolvedRegionIR(
                region_source_id=unresolved_id,
                page_id=page_id,
                bbox=_full_page_bbox(dimensions),
                reason_code="empty_text_document",
                producer=_PRODUCER,
                evidence_refs=(render_ref,),
            ),
        )
    page = PageIR(
        page_id=page_id,
        physical_index=0,
        dimensions=dimensions,
        blocks=ordered_blocks,
        native_objects=tuple(objects),
        observed_order_signals=signals,
        reading_order_hypotheses=reading,
        role_hypotheses=(
            _role_hypothesis(
                source_hash=source_hash,
                page_id=page_id,
                physical_index=0,
                text=text,
                title=title,
                content_block_count=len(blocks),
            ),
        ),
        render_ref=render_ref,
    )
    return (page,), tuple(outline), unresolved


def parse_source(path: Path) -> SourceObservationIR:
    """Parse one source file into the immutable vNext observation contract."""

    source_path = path.resolve()
    if not source_path.is_file():
        raise FileNotFoundError(source_path)
    suffix = source_path.suffix.casefold()
    if suffix not in SUPPORTED_SUFFIXES:
        supported = ", ".join(sorted(SUPPORTED_SUFFIXES))
        raise ValueError(
            f"unsupported source type {suffix or '<none>'}; expected {supported}"
        )
    source_hash = _source_hash(source_path)
    document_id = _source_id(
        "document",
        source_hash=source_hash,
        locator={"revision": 1},
    )
    if suffix == ".pdf":
        pages, outline, unresolved = _parse_pdf(
            source_path,
            source_hash=source_hash,
        )
    elif suffix == ".pptx":
        pages, outline, unresolved = _parse_pptx(
            source_path,
            source_hash=source_hash,
        )
    elif suffix == ".docx":
        pages, outline, unresolved = _parse_docx(
            source_path,
            source_hash=source_hash,
        )
    else:
        pages, outline, unresolved = _parse_text(
            source_path,
            source_hash=source_hash,
        )
    return SourceObservationIR(
        document_id=document_id,
        source_hash=source_hash,
        parser_manifest=_parser_manifest(suffix),
        pages=pages,
        outline_entries=outline,
        unresolved_regions=unresolved,
    )
